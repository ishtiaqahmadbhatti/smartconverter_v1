import os
import json
import csv
import tempfile
from typing import List, Dict, Any, Optional, Union
from pathlib import Path
import fitz  # PyMuPDF
import pandas as pd
from PIL import Image
import markdown
from docx import Document
from pptx import Presentation
import openpyxl
from openpyxl import Workbook
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
try:
    from weasyprint import HTML, CSS
    from weasyprint.text.fonts import FontConfiguration
    WEASYPRINT_AVAILABLE = True
except (ImportError, OSError):
    WEASYPRINT_AVAILABLE = False
    HTML = None
    CSS = None
    FontConfiguration = None
try:
    import cairosvg
    CAIROSVG_AVAILABLE = True
except ImportError:
    CAIROSVG_AVAILABLE = False
from io import BytesIO
import base64
import re
import zipfile
from datetime import datetime
from PyPDF2 import PdfReader, PdfWriter
import shutil
import subprocess
from app.core.config import settings
from app.core.exceptions import FileProcessingError
from PyPDF2 import PdfMerger


class PDFConversionService:
    """Service for comprehensive PDF conversion operations."""
    
    # Supported input formats
    SUPPORTED_INPUT_FORMATS = [
        "PDF", "HTML", "DOCX", "PPTX", "OXPS", "JPG", "JPEG", "PNG", 
        "MD", "XLSX", "XLS", "ODS", "CSV", "TXT"
    ]
    
    # Supported output formats
    SUPPORTED_OUTPUT_FORMATS = [
        "PDF", "JSON", "MD", "CSV", "XLSX", "XPS", "JPG", "PNG", 
        "TIFF", "SVG", "HTML", "TXT", "DOCX"
    ]
    
    @staticmethod
    def merge_pdfs(input_paths: List[str], output_path: str) -> str:
        """Merge multiple PDF files into a single PDF."""
        if not input_paths:
            raise FileProcessingError("No PDF files provided for merging.")

        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            merger = PdfMerger()

            for path in input_paths:
                if not os.path.exists(path):
                    raise FileProcessingError(f"Input file not found: {path}")
                with open(path, "rb") as pdf_file:
                    merger.append(pdf_file)

            with open(output_path, "wb") as merged_file:
                merger.write(merged_file)
            merger.close()

            return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(f"Error merging PDFs: {str(e)}")
    
    @staticmethod
    def pdf_to_json(pdf_path: str, output_path: str) -> str:
        """Convert PDF to JSON format with structured data extraction."""
        try:
            doc = fitz.open(pdf_path)
            pages_data = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text
                text = page.get_text()
                
                # Extract images
                images = []
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                        images.append({
                            "index": img_index,
                            "format": "png",
                            "data": base64.b64encode(img_data).decode()
                        })
                    pix = None
                
                # Extract annotations
                annotations = []
                for annot in page.annots():
                    annotations.append({
                        "type": annot.type[1],
                        "content": annot.content,
                        "rect": list(annot.rect)
                    })
                
                # Extract tables (basic)
                tables = []
                # This is a simplified table extraction
                # In production, you'd want to use more sophisticated table detection
                
                page_data = {
                    "page_number": page_num + 1,
                    "text": text,
                    "images": images,
                    "annotations": annotations,
                    "tables": tables,
                    "dimensions": {
                        "width": page.rect.width,
                        "height": page.rect.height
                    }
                }
                pages_data.append(page_data)
            
            doc.close()
            
            # Create structured JSON
            pdf_data = {
                "document_info": {
                    "total_pages": len(pages_data),
                    "title": doc.metadata.get("title", ""),
                    "author": doc.metadata.get("author", ""),
                    "subject": doc.metadata.get("subject", ""),
                    "creator": doc.metadata.get("creator", ""),
                    "producer": doc.metadata.get("producer", ""),
                    "creation_date": doc.metadata.get("creationDate", ""),
                    "modification_date": doc.metadata.get("modDate", "")
                },
                "pages": pages_data
            }
            
            # Save JSON
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(pdf_data, f, indent=2, ensure_ascii=False)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to JSON: {str(e)}")
    
    @staticmethod
    def pdf_to_markdown(pdf_path: str, output_path: str) -> str:
        """Convert PDF to Markdown format."""
        try:
            doc = fitz.open(pdf_path)
            markdown_content = []
            
            # Add document metadata
            metadata = doc.metadata
            if metadata.get("title"):
                markdown_content.append(f"# {metadata['title']}\n")
            if metadata.get("author"):
                markdown_content.append(f"**Author:** {metadata['author']}\n")
            if metadata.get("subject"):
                markdown_content.append(f"**Subject:** {metadata['subject']}\n")
            markdown_content.append("---\n")
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    markdown_content.append(f"## Page {page_num + 1}\n")
                    markdown_content.append(text)
                    markdown_content.append("\n---\n")
            
            doc.close()
            
            # Save markdown
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(markdown_content))
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to Markdown: {str(e)}")
    
    @staticmethod
    def pdf_to_csv(pdf_path: str, output_path: str) -> str:
        """Convert PDF to CSV format (extract tabular data)."""
        try:
            doc = fitz.open(pdf_path)
            all_tables = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract tables using PyMuPDF's table detection
                tables = page.find_tables()
                for table in tables:
                    table_data = table.extract()
                    if table_data:
                        all_tables.extend(table_data)
            
            doc.close()
            
            # Save as CSV
            if all_tables:
                with open(output_path, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    writer.writerows(all_tables)
            else:
                # If no tables found, create a simple text-based CSV
                with open(output_path, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    writer.writerow(['Page', 'Content'])
                    doc = fitz.open(pdf_path)
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        text = page.get_text()
                        if text.strip():
                            writer.writerow([page_num + 1, text.strip()])
                    doc.close()
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to CSV: {str(e)}")
    
    @staticmethod
    def pdf_to_excel(pdf_path: str, output_path: str) -> str:
        """Convert PDF to Excel format."""
        try:
            doc = fitz.open(pdf_path)
            workbook = Workbook()
            workbook.remove(workbook.active)  # Remove default sheet
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Create a new sheet for each page
                ws = workbook.create_sheet(title=f"Page_{page_num + 1}")
                
                # Extract tables
                tables = list(page.find_tables())  # Convert to list to check if empty
                row_offset = 1
                has_table_data = False
                
                # Extract and add table data
                for table in tables:
                    try:
                        table_data = table.extract()
                        if table_data and len(table_data) > 0:
                            has_table_data = True
                            for row_idx, row in enumerate(table_data):
                                for col_idx, cell in enumerate(row):
                                    if cell:  # Only add non-empty cells
                                        ws.cell(row=row_offset + row_idx, column=col_idx + 1, value=str(cell))
                            row_offset += len(table_data) + 2  # Add spacing between tables
                    except Exception as e:
                        # If table extraction fails, continue with text extraction
                        pass
                
                # Extract text content
                text = page.get_text()
                if text and text.strip():
                    # If we have table data, add text after tables, otherwise start from row 1
                    if not has_table_data:
                        row_offset = 1
                    
                    lines = text.split('\n')
                    # Filter out empty lines and limit to reasonable number
                    non_empty_lines = [line.strip() for line in lines if line.strip()][:500]
                    
                    if non_empty_lines:
                        # Add header if we have both tables and text
                        if has_table_data:
                            ws.cell(row=row_offset, column=1, value="Text Content:")
                            row_offset += 1
                        
                        for row_idx, line in enumerate(non_empty_lines):
                            ws.cell(row=row_offset + row_idx, column=1, value=line)
            
            doc.close()
            workbook.save(output_path)
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to Excel: {str(e)}")
    
    @staticmethod
    def html_to_pdf(html_path: str, output_path: str) -> str:
        """Convert HTML to PDF."""
        try:
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            if not WEASYPRINT_AVAILABLE:
                # Fallback to reportlab for basic HTML to PDF conversion
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.platypus import SimpleDocTemplate, Paragraph
                from reportlab.lib.pagesizes import letter
                
                # Simple HTML to text conversion for reportlab
                import re
                text_content = re.sub(r'<[^>]+>', '', html_content)
                
                doc = SimpleDocTemplate(output_path, pagesize=letter)
                styles = getSampleStyleSheet()
                story = []
                
                # Split content into paragraphs
                paragraphs = text_content.split('\n\n')
                for para in paragraphs:
                    if para.strip():
                        story.append(Paragraph(para.strip(), styles['Normal']))
                
                doc.build(story)
            else:
                # Use WeasyPrint for proper HTML to PDF conversion
                font_config = FontConfiguration()
                HTML(string=html_content).write_pdf(output_path, font_config=font_config)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting HTML to PDF: {str(e)}")
    
    @staticmethod
    def word_to_pdf(docx_path: str, output_path: str) -> str:
        """Convert Word document to PDF."""
        try:
            # This is a simplified conversion
            # In production, you might want to use python-docx2pdf or other libraries
            doc = Document(docx_path)
            
            # Create PDF using reportlab
            doc_pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    p = Paragraph(paragraph.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            doc_pdf.build(story)
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting Word to PDF: {str(e)}")
    
    @staticmethod
    def powerpoint_to_pdf(pptx_path: str, output_path: str) -> str:
        """Convert PowerPoint to PDF."""
        try:
            prs = Presentation(pptx_path)
            
            # Create PDF using reportlab
            doc_pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for slide_num, slide in enumerate(prs.slides):
                story.append(Paragraph(f"Slide {slide_num + 1}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        story.append(Paragraph(shape.text, styles['Normal']))
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 20))
            
            doc_pdf.build(story)
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PowerPoint to PDF: {str(e)}")
    
    @staticmethod
    def oxps_to_pdf(oxps_path: str, output_path: str) -> str:
        """Convert OXPS to PDF."""
        try:
            # Use PyMuPDF (fitz) to convert OXPS to PDF
            doc = fitz.open(oxps_path)
            pdf_bytes = doc.convert_to_pdf()
            with open(output_path, "wb") as f:
                f.write(pdf_bytes)
            doc.close()
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting OXPS to PDF: {str(e)}")
    
    @staticmethod
    def image_to_pdf(image_path: str, output_path: str) -> str:
        """Convert image (JPG/PNG) to PDF."""
        try:
            img = Image.open(image_path)
            
            # Convert to RGB if necessary
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Save as PDF
            img.save(output_path, 'PDF', resolution=300.0)
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting image to PDF: {str(e)}")
    
    @staticmethod
    def markdown_to_pdf(md_path: str, output_path: str) -> str:
        """Convert Markdown to PDF with proper handling of emojis, headings, and tables."""
        try:
            with open(md_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # Convert markdown to HTML with extensions for tables and other features
            try:
                # Try to use markdown extensions if available
                import markdown.extensions.tables
                import markdown.extensions.fenced_code
                md_extensions = ['tables', 'fenced_code', 'nl2br']
                html_content = markdown.markdown(md_content, extensions=md_extensions)
            except:
                # Fallback to basic markdown if extensions not available
                html_content = markdown.markdown(md_content)
            
            
            # custom flag to track if we need fallback
            weasyprint_success = False

            # Try WeasyPrint FIRST
            if WEASYPRINT_AVAILABLE:
                try:
                    # Use WeasyPrint for proper HTML to PDF conversion with CSS styling
                    css_content = """
                    @page {
                        size: A4;
                        margin: 2cm;
                    }
                    body {
                        font-family: 'DejaVu Sans', Arial, sans-serif;
                        font-size: 11pt;
                        line-height: 1.6;
                        color: #333;
                    }
                    h1 {
                        font-size: 24pt;
                        color: #1a1a1a;
                        margin-top: 20pt;
                        margin-bottom: 12pt;
                        border-bottom: 2px solid #ddd;
                        padding-bottom: 8pt;
                    }
                    h2 {
                        font-size: 20pt;
                        color: #2a2a2a;
                        margin-top: 16pt;
                        margin-bottom: 10pt;
                        border-bottom: 1px solid #eee;
                        padding-bottom: 6pt;
                    }
                    h3 {
                        font-size: 16pt;
                        color: #3a3a3a;
                        margin-top: 12pt;
                        margin-bottom: 8pt;
                    }
                    h4 {
                        font-size: 14pt;
                        color: #4a4a4a;
                        margin-top: 10pt;
                        margin-bottom: 6pt;
                    }
                    table {
                        border-collapse: collapse;
                        width: 100%;
                        margin: 12pt 0;
                    }
                    th {
                        background-color: #f0f0f0;
                        color: #000;
                        font-weight: bold;
                        padding: 8pt;
                        border: 1px solid #ddd;
                    }
                    td {
                        padding: 6pt;
                        border: 1px solid #ddd;
                    }
                    tr:nth-child(even) {
                        background-color: #f9f9f9;
                    }
                    p {
                        margin: 8pt 0;
                    }
                    ul, ol {
                        margin: 8pt 0;
                        padding-left: 24pt;
                    }
                    li {
                        margin: 4pt 0;
                    }
                    code {
                        background-color: #f4f4f4;
                        padding: 2pt 4pt;
                        border-radius: 3pt;
                        font-family: 'Courier New', monospace;
                    }
                    pre {
                        background-color: #f5f5f5;
                        border: 1px solid #ddd;
                        border-radius: 4pt;
                        padding: 12pt;
                        margin: 12pt 0;
                        overflow-x: auto;
                        font-family: 'Courier New', 'Courier', monospace;
                        font-size: 8pt;
                        line-height: 1.4;
                        white-space: pre;
                        word-wrap: normal;
                    }
                    pre code {
                        background-color: transparent;
                        padding: 0;
                        border-radius: 0;
                        font-size: inherit;
                    }
                    """
                    
                    font_config = FontConfiguration()
                    HTML(string=html_content).write_pdf(
                        output_path,
                        font_config=font_config,
                        stylesheets=[CSS(string=css_content)]
                    )
                    
                    weasyprint_success = True
                    return output_path
                except Exception as e:
                    print(f"WeasyPrint failed: {e}. Falling back to ReportLab.")
                    weasyprint_success = False

            if not weasyprint_success:
                # Fallback to reportlab with proper HTML parsing using BeautifulSoup
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
                from reportlab.lib.pagesizes import A4
                from reportlab.lib import colors
                from reportlab.lib.units import inch
                from bs4 import BeautifulSoup
                import re
                
                # Parse HTML with BeautifulSoup for better handling of nested tags
                soup = BeautifulSoup(html_content, 'html.parser')
                story = []
                styles = getSampleStyleSheet()
                
                # Register Unicode fonts for emoji support
                unicode_font_name = 'Helvetica'  # Default fallback
                try:
                    from reportlab.pdfbase import pdfmetrics
                    from reportlab.pdfbase.ttfonts import TTFont
                    import sys
                    import os
                    
                    # Try to register fonts with better Unicode/Emoji support
                    # Try multiple font paths that might support emojis
                    try:
                        # Common system font paths (prioritize fonts with emoji support)
                        font_paths = [
                            # Windows fonts (better emoji support)
                            'C:/Windows/Fonts/seguiemj.ttf',  # Segoe UI Emoji
                            'C:/Windows/Fonts/segmdl2.ttf',   # Segoe MDL2 Assets
                            'C:/Windows/Fonts/arial.ttf',     # Arial (basic Unicode)
                            'C:/Windows/Fonts/calibri.ttf',  # Calibri
                            # Linux fonts
                            '/usr/share/fonts/truetype/noto/NotoColorEmoji.ttf',
                            '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
                            '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
                            # macOS fonts
                            '/System/Library/Fonts/Supplemental/Apple Color Emoji.ttc',
                            '/System/Library/Fonts/Helvetica.ttc',
                        ]
                        
                        # Try to find and register a Unicode font
                        for font_path in font_paths:
                            if os.path.exists(font_path):
                                try:
                                    pdfmetrics.registerFont(TTFont('UnicodeFont', font_path))
                                    unicode_font_name = 'UnicodeFont'
                                    break
                                except Exception as e:
                                    # Continue trying other fonts
                                    continue
                    except:
                        pass
                except:
                    pass  # Keep default Helvetica
                
                # Define heading styles with Unicode font support
                heading_styles = {
                    'h1': ParagraphStyle(
                        'CustomH1',
                        parent=styles['Heading1'],
                        fontSize=24,
                        spaceAfter=12,
                        textColor=colors.HexColor('#1a1a1a'),
                        fontName=unicode_font_name
                    ),
                    'h2': ParagraphStyle(
                        'CustomH2',
                        parent=styles['Heading2'],
                        fontSize=20,
                        spaceAfter=10,
                        textColor=colors.HexColor('#2a2a2a'),
                        fontName=unicode_font_name
                    ),
                    'h3': ParagraphStyle(
                        'CustomH3',
                        parent=styles['Heading3'],
                        fontSize=16,
                        spaceAfter=8,
                        textColor=colors.HexColor('#3a3a3a'),
                        fontName=unicode_font_name
                    ),
                    'h4': ParagraphStyle(
                        'CustomH4',
                        parent=styles['Heading4'],
                        fontSize=14,
                        spaceAfter=6,
                        fontName=unicode_font_name
                    ),
                    'h5': ParagraphStyle(
                        'CustomH5',
                        parent=styles['Heading5'],
                        fontSize=12,
                        spaceAfter=4,
                        fontName=unicode_font_name
                    ),
                    'h6': ParagraphStyle(
                        'CustomH6',
                        parent=styles['Heading6'],
                        fontSize=11,
                        spaceAfter=4,
                        fontName=unicode_font_name
                    ),
                }
                
                # Normal style with Unicode font
                normal_style = ParagraphStyle(
                    'UnicodeNormal',
                    parent=styles['Normal'],
                    fontName=unicode_font_name,
                    encoding='utf-8'
                )
                
                def get_text_with_emojis(element):
                    """Extract text from element preserving emojis and special characters."""
                    if element is None:
                        return ""
                    # Get all text including emojis - preserve Unicode characters
                    text = element.get_text(separator=' ', strip=False)
                    # Clean up extra whitespace but preserve structure and emojis
                    text = re.sub(r'\s+', ' ', text).strip()
                    # Ensure text is properly encoded as UTF-8
                    if isinstance(text, bytes):
                        text = text.decode('utf-8', errors='ignore')
                    return text
                
                def create_paragraph_with_emojis(text, style):
                    """Create a paragraph that properly handles emojis."""
                    if not text:
                        return None
                    try:
                        # Ensure text is UTF-8 encoded
                        if isinstance(text, bytes):
                            text = text.decode('utf-8', errors='ignore')
                        # Escape XML special characters but preserve emojis
                        from xml.sax.saxutils import escape
                        # Don't escape emojis - they should be preserved as-is
                        # Only escape XML special chars
                        text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        # Create paragraph with UTF-8 encoding
                        return Paragraph(text, style)
                    except Exception as e:
                        # Fallback: try without escaping
                        try:
                            return Paragraph(str(text), style)
                        except:
                            # Last resort: remove problematic characters but keep emojis
                            safe_text = ''.join(c if ord(c) < 0x10000 or c.isprintable() else '?' for c in str(text))
                            return Paragraph(safe_text, style)
                
                def process_element(element):
                    """Recursively process HTML elements and add to story."""
                    if element is None:
                        return
                    
                    tag_name = element.name if hasattr(element, 'name') else None
                    
                    if tag_name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                        text = get_text_with_emojis(element)
                        if text:
                            style = heading_styles.get(tag_name, styles['Heading1'])
                            para = create_paragraph_with_emojis(text, style)
                            if para:
                                story.append(para)
                                story.append(Spacer(1, 6))
                    
                    elif tag_name == 'p':
                        text = get_text_with_emojis(element)
                        if text:
                            para = create_paragraph_with_emojis(text, normal_style)
                            if para:
                                story.append(para)
                                story.append(Spacer(1, 8))
                    
                    elif tag_name in ['ul', 'ol']:
                        for li in element.find_all('li', recursive=False):
                            text = get_text_with_emojis(li)
                            if text:
                                bullet = "•" if tag_name == 'ul' else "1."
                                para = create_paragraph_with_emojis(f"{bullet} {text}", normal_style)
                                if para:
                                    story.append(para)
                                    story.append(Spacer(1, 4))
                        story.append(Spacer(1, 6))
                    
                    elif tag_name == 'table':
                        # Extract table data
                        table_data = []
                        rows = element.find_all('tr')
                        for row in rows:
                            cells = row.find_all(['td', 'th'])
                            row_data = [get_text_with_emojis(cell) for cell in cells]
                            if row_data:
                                table_data.append(row_data)
                        
                        if table_data:
                            # Convert table cells to Paragraphs to support emojis
                            formatted_table_data = []
                            for row_idx, row in enumerate(table_data):
                                formatted_row = []
                                for cell_text in row:
                                    if cell_text:
                                        # Create paragraph for each cell to support emojis
                                        cell_para = create_paragraph_with_emojis(
                                            cell_text, 
                                            normal_style if row_idx > 0 else heading_styles.get('h4', normal_style)
                                        )
                                        formatted_row.append(cell_para if cell_para else cell_text)
                                    else:
                                        formatted_row.append('')
                                formatted_table_data.append(formatted_row)
                            
                            # Create table with proper styling
                            pdf_table = Table(formatted_table_data)
                            pdf_table.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                                ('FONTNAME', (0, 0), (-1, 0), unicode_font_name),
                                ('FONTSIZE', (0, 0), (-1, 0), 11),
                                ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
                                ('TOPPADDING', (0, 0), (-1, 0), 10),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                                ('FONTSIZE', (0, 1), (-1, -1), 10),
                                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
                                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                            ]))
                            story.append(pdf_table)
                            story.append(Spacer(1, 12))
                    
                    elif tag_name == 'hr':
                        story.append(Spacer(1, 12))
                        # Add a line separator
                        from reportlab.platypus import HRFlowable
                        story.append(HRFlowable(width="100%", thickness=1, lineCap='round', color=colors.grey))
                        story.append(Spacer(1, 12))
                    
                    elif tag_name == 'br':
                        story.append(Spacer(1, 6))
                    
                    elif tag_name == 'pre':
                        # Handle pre-formatted code blocks (from triple backticks)
                        # Extract all text preserving line breaks and formatting
                        code_text = ""
                        if element.find('code'):
                            # If there's a code tag inside pre, get its text
                            code_text = element.find('code').get_text(separator='\n', strip=False)
                        else:
                            # Otherwise get all text from pre tag
                            code_text = element.get_text(separator='\n', strip=False)
                        
                        if code_text:
                            # Preserve line breaks and formatting
                            code_text = code_text.rstrip()  # Remove trailing newlines
                            
                            # Create a code block style with monospace font
                            pre_code_style = ParagraphStyle(
                                'PreCode',
                                parent=styles['Normal'],
                                fontName='Courier',
                                fontSize=8,
                                leading=10,
                                encoding='utf-8',
                                leftIndent=0,
                                rightIndent=0
                            )
                            
                            # Split by lines and create paragraphs to preserve formatting
                            lines = code_text.split('\n')
                            code_paragraphs = []
                            for line in lines:
                                # Preserve all lines including empty ones for proper formatting
                                # Escape XML special chars but preserve structure
                                safe_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                # Preserve spaces at the beginning (for indentation)
                                para = Paragraph(safe_line, pre_code_style)
                                code_paragraphs.append(para)
                            
                            # Wrap code block in a Table to create background and border effect
                            # Each line becomes a cell in a single-column table
                            code_table_data = [[para] for para in code_paragraphs]
                            code_table = Table(code_table_data, colWidths=[None])  # Auto width
                            
                            # Style the table to look like a code block
                            code_table.setStyle(TableStyle([
                                # Background color for all cells
                                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#f5f5f5')),
                                # Border around the entire block
                                ('BOX', (0, 0), (-1, -1), 1, colors.HexColor('#ddd')),
                                # Padding inside cells
                                ('LEFTPADDING', (0, 0), (-1, -1), 12),
                                ('RIGHTPADDING', (0, 0), (-1, -1), 12),
                                ('TOPPADDING', (0, 0), (-1, -1), 6),
                                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                                # Alignment
                                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                            ]))
                            
                            story.append(Spacer(1, 8))
                            story.append(code_table)
                            story.append(Spacer(1, 8))
                    
                    elif tag_name == 'code':
                        # Handle inline code (single backticks)
                        text = get_text_with_emojis(element)
                        if text:
                            code_style = ParagraphStyle(
                                'Code',
                                parent=styles['Normal'],
                                fontName='Courier',
                                fontSize=9,
                                backColor=colors.lightgrey,
                                leftIndent=12,
                                rightIndent=12,
                                spaceBefore=4,
                                spaceAfter=4,
                                encoding='utf-8'
                            )
                            para = create_paragraph_with_emojis(text, code_style)
                            if para:
                                story.append(para)
                                story.append(Spacer(1, 6))
                    
                    elif tag_name == 'blockquote':
                        text = get_text_with_emojis(element)
                        if text:
                            quote_style = ParagraphStyle(
                                'Quote',
                                parent=normal_style,
                                leftIndent=20,
                                rightIndent=20,
                                fontStyle='Italic',
                                textColor=colors.HexColor('#555555')
                            )
                            para = create_paragraph_with_emojis(f'"{text}"', quote_style)
                            if para:
                                story.append(para)
                                story.append(Spacer(1, 8))
                    
                    else:
                        # Process children recursively
                        if hasattr(element, 'children'):
                            for child in element.children:
                                if hasattr(child, 'name'):
                                    process_element(child)
                                elif isinstance(child, str) and child.strip():
                                    # Handle text nodes
                                    text = child.strip()
                                    if text:
                                        para = create_paragraph_with_emojis(text, normal_style)
                                        if para:
                                            story.append(para)
                                            story.append(Spacer(1, 6))
                
                # Process all top-level elements
                for element in soup.children:
                    if hasattr(element, 'name'):
                        process_element(element)
                
                # If no content was added, add body content
                if not story:
                    body = soup.find('body')
                    if body:
                        for element in body.children:
                            if hasattr(element, 'name'):
                                process_element(element)
                    else:
                        # Fallback: process all elements
                        for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'table']):
                            process_element(element)
                
                # Create PDF
                doc = SimpleDocTemplate(
                    output_path,
                    pagesize=A4,
                    rightMargin=72,
                    leftMargin=72,
                    topMargin=72,
                    bottomMargin=72,
                    encoding='utf-8'
                )
                doc.build(story)

            
        except Exception as e:
            raise FileProcessingError(f"Error converting Markdown to PDF: {str(e)}")
    
    @staticmethod
    def excel_to_pdf(xlsx_path: str, output_path: str) -> str:
        """Convert Excel to PDF."""
        try:
            workbook = openpyxl.load_workbook(xlsx_path)
            
            # Create PDF using reportlab
            doc_pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_name in workbook.sheetnames:
                ws = workbook[sheet_name]
                story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                # Convert sheet data to table
                data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append([str(cell) if cell is not None else "" for cell in row])
                
                if data:
                    from reportlab.platypus import Table
                    table = Table(data)
                    story.append(table)
                    story.append(Spacer(1, 20))
            
            doc_pdf.build(story)
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting Excel to PDF: {str(e)}")
    
    @staticmethod
    def excel_to_xps(xlsx_path: str, output_path: str) -> str:
        """Convert Excel to XPS."""
        try:
            # XPS conversion is complex and requires specialized libraries
            # This is a placeholder implementation
            raise FileProcessingError("Excel to XPS conversion requires specialized libraries not available in standard Python packages")
            
        except Exception as e:
            raise FileProcessingError(f"Error converting Excel to XPS: {str(e)}")
    
    @staticmethod
    def ods_to_pdf(ods_path: str, output_path: str) -> str:
        """Convert OpenOffice Calc ODS to PDF."""
        try:
            # ODS conversion requires specialized libraries
            # This is a placeholder implementation
            raise FileProcessingError("ODS to PDF conversion requires specialized libraries not available in standard Python packages")
            
        except Exception as e:
            raise FileProcessingError(f"Error converting ODS to PDF: {str(e)}")
    
    @staticmethod
    def pdf_to_csv_extract(pdf_path: str, output_path: str) -> str:
        """Extract tabular data from PDF to CSV."""
        return PDFConversionService.pdf_to_csv(pdf_path, output_path)
    
    @staticmethod
    def pdf_to_excel_extract(pdf_path: str, output_path: str) -> str:
        """Extract tabular data from PDF to Excel."""
        return PDFConversionService.pdf_to_excel(pdf_path, output_path)
    
    @staticmethod
    def pdf_to_word_extract(pdf_path: str, output_path: str) -> str:
        """Convert PDF to Word document."""
        try:
            doc = fitz.open(pdf_path)
            word_doc = Document()
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    # Add page heading
                    word_doc.add_heading(f'Page {page_num + 1}', level=2)
                    
                    # Add text content
                    paragraphs = text.split('\n')
                    for para in paragraphs:
                        if para.strip():
                            word_doc.add_paragraph(para.strip())
                    
                    # Add page break (except for last page)
                    if page_num < len(doc) - 1:
                        word_doc.add_page_break()
            
            doc.close()
            word_doc.save(output_path)
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to Word: {str(e)}")
    
    @staticmethod
    def pdf_to_image(pdf_path: str, output_dir: str, format: str = "jpg") -> List[str]:
        """Convert PDF pages to images."""
        try:
            doc = fitz.open(pdf_path)
            output_files = []
            target_format = (format or "jpg").lower()
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()
                
                output_file = os.path.join(output_dir, f"page_{page_num + 1}.{target_format}")
                
                if target_format in {"tiff", "tif"}:
                    mode = "RGBA" if pix.alpha else "RGB"
                    image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
                    image.save(output_file, format="TIFF")
                else:
                    pix.save(output_file)
                output_files.append(output_file)
            
            doc.close()
            return output_files
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to images: {str(e)}")
    
    @staticmethod
    def pdf_to_tiff(pdf_path: str, output_dir: str) -> List[str]:
        """Convert PDF pages to TIFF images."""
        return PDFConversionService.pdf_to_image(pdf_path, output_dir, "tiff")
    
    @staticmethod
    def pdf_to_svg(pdf_path: str, output_dir: str) -> List[str]:
        """Convert PDF pages to SVG."""
        try:
            doc = fitz.open(pdf_path)
            output_files = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                svg_string = page.get_svg_image()
                
                output_file = os.path.join(output_dir, f"page_{page_num + 1}.svg")
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(svg_string)
                output_files.append(output_file)
            
            doc.close()
            return output_files
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to SVG: {str(e)}")
    
    @staticmethod
    def pdf_to_html(pdf_path: str, output_path: str) -> str:
        """Convert PDF to HTML."""
        try:
            doc = fitz.open(pdf_path)
            html_content = []
            
            # Add HTML header
            html_content.append("<!DOCTYPE html>")
            html_content.append("<html><head><title>PDF Content</title></head><body>")
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    html_content.append(f"<div class='page' id='page-{page_num + 1}'>")
                    html_content.append(f"<h2>Page {page_num + 1}</h2>")
                    html_content.append(f"<p>{text.replace(chr(10), '<br>')}</p>")
                    html_content.append("</div>")
            
            html_content.append("</body></html>")
            
            doc.close()
            
            # Save HTML
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(html_content))
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to HTML: {str(e)}")
    
    @staticmethod
    def pdf_to_text(pdf_path: str, output_path: str) -> str:
        """Convert PDF to plain text."""
        try:
            doc = fitz.open(pdf_path)
            text_content = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    text_content.append(f"--- Page {page_num + 1} ---")
                    text_content.append(text)
                    text_content.append("")
            
            doc.close()
            
            # Save text
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(text_content))
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Error converting PDF to text: {str(e)}")

    @staticmethod
    def crop_pdf(input_path: str, output_path: str, crop_box: Dict[str, int]) -> str:
        try:
            x = int(crop_box.get("x", 0))
            y = int(crop_box.get("y", 0))
            w = int(crop_box.get("width", 0))
            h = int(crop_box.get("height", 0))
            if w <= 0 or h <= 0:
                raise FileProcessingError("Invalid crop size")

            try:
                doc = fitz.open(input_path)
                for page in doc:
                    r = page.rect
                    x1 = max(r.x0, r.x0 + x)
                    y1 = max(r.y0, r.y0 + y)
                    x2 = min(r.x1, x1 + w)
                    y2 = min(r.y1, y1 + h)
                    clip = fitz.Rect(x1, y1, x2, y2)
                    pm = page.get_pixmap(clip=clip)
                    img = Image.frombytes("RGBA" if pm.alpha else "RGB", [pm.width, pm.height], pm.samples)
                    mem = BytesIO()
                    img.save(mem, format="PNG")
                    mem.seek(0)
                    # Replace page content with the cropped image
                    new_page = doc.new_page(-1, width=pm.width, height=pm.height)
                    rect = fitz.Rect(0, 0, pm.width, pm.height)
                    new_page.insert_image(rect, stream=mem.read())
                # Remove original pages
                for _ in range(len(doc) // 2):
                    doc.delete_page(0)
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                doc.save(output_path)
                doc.close()
                return output_path
            except Exception:
                reader = PdfReader(input_path)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, "wb") as f:
                    writer.write(f)
                return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def protect_pdf(input_path: str, output_path: str, password: str) -> str:
        try:
            reader = PdfReader(input_path)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            if not password:
                raise FileProcessingError("Password cannot be empty")
            writer.encrypt(password)
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path, "wb") as f:
                writer.write(f)
            return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def unlock_pdf(input_path: str, output_path: str, password: str) -> str:
        try:
            try:
                reader = PdfReader(input_path)
                if reader.is_encrypted:
                    reader.decrypt(password)
            except Exception:
                reader = PdfReader(input_path)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path, "wb") as f:
                writer.write(f)
            return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def repair_pdf(input_path: str, output_path: str) -> str:
        try:
            try:
                doc = fitz.open(input_path)
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                doc.save(output_path, clean=True, deflate=True, garbage=2)
                doc.close()
                return output_path
            except Exception:
                reader = PdfReader(input_path)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, "wb") as f:
                    writer.write(f)
                return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def compare_pdfs(input_path1: str, input_path2: str, output_path: str) -> Dict[str, Any]:
        try:
            doc1 = fitz.open(input_path1)
            doc2 = fitz.open(input_path2)
            pages1 = len(doc1)
            pages2 = len(doc2)
            max_pages = max(pages1, pages2)
            diffs: List[Dict[str, Any]] = []
            for i in range(max_pages):
                t1 = doc1.load_page(i).get_text() if i < pages1 else ""
                t2 = doc2.load_page(i).get_text() if i < pages2 else ""
                if (t1 or "").strip() != (t2 or "").strip():
                    diffs.append({"page": i + 1, "difference": True})
            doc1.close()
            doc2.close()
            summary = {
                "pages_doc1": pages1,
                "pages_doc2": pages2,
                "differences_count": len(diffs),
                "differences": diffs,
            }
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path, "w", encoding="utf-8") as f:
                import json as _json
                _json.dump(summary, f, indent=2, ensure_ascii=False)
            return summary
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def get_pdf_metadata(input_path: str) -> Dict[str, Any]:
        try:
            doc = fitz.open(input_path)
            meta = {
                "page_count": len(doc),
                "metadata": doc.metadata,
            }
            doc.close()
            return meta
        except Exception as e:
            raise FileProcessingError(str(e))
    
    @staticmethod
    def get_supported_formats() -> Dict[str, List[str]]:
        """Get supported input and output formats."""
        return {
            "input_formats": PDFConversionService.SUPPORTED_INPUT_FORMATS,
            "output_formats": PDFConversionService.SUPPORTED_OUTPUT_FORMATS
        }
    
    @staticmethod
    def cleanup_temp_files(file_path: str):
        """Clean up temporary files."""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception:
            pass  # Ignore cleanup errors

    @staticmethod
    def compress_pdf(
        input_path: str,
        output_path: str,
        compression_level: str = "medium",
        target_reduction_pct: Optional[int] = None,
        max_image_dpi: Optional[int] = None,
    ) -> str:
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            level_key = (compression_level or "medium").lower()
            gs_settings_order = {
                "low": ["/printer", "/ebook"],
                "medium": ["/ebook", "/screen"],
                "high": ["/screen"],
            }
            original_size = os.path.getsize(input_path) if os.path.exists(input_path) else None
            desired_size = None
            if original_size and isinstance(target_reduction_pct, int):
                if target_reduction_pct < 0:
                    target_reduction_pct = 0
                if target_reduction_pct > 100:
                    target_reduction_pct = 100
                desired_size = int(original_size * (1 - (target_reduction_pct / 100.0)))

            gs_exe = shutil.which("gswin64c") or shutil.which("gswin32c") or shutil.which("gs")
            used_method = None

            def _run_gs(setting: str, dpi: Optional[int]) -> Optional[int]:
                try:
                    cmd = [
                        gs_exe,
                        "-sDEVICE=pdfwrite",
                        "-dCompatibilityLevel=1.4",
                        f"-dPDFSETTINGS={setting}",
                        "-dNOPAUSE",
                        "-dQUIET",
                        "-dBATCH",
                    ]
                    if dpi:
                        cmd += [
                            "-dDownsampleColorImages=true",
                            f"-dColorImageResolution={dpi}",
                            "-dColorImageDownsampleType=/Average",
                            "-dDownsampleGrayImages=true",
                            f"-dGrayImageResolution={dpi}",
                            "-dGrayImageDownsampleType=/Average",
                            "-dDownsampleMonoImages=true",
                            f"-dMonoImageResolution={dpi}",
                            "-dMonoImageDownsampleType=/Subsample",
                        ]
                    cmd += [f"-sOutputFile={output_path}", input_path]
                    subprocess.run(cmd, check=True)
                    if os.path.exists(output_path):
                        return os.path.getsize(output_path)
                    return None
                except Exception:
                    return None

            if gs_exe:
                dpi_candidates = []
                if isinstance(max_image_dpi, int) and max_image_dpi > 0:
                    dpi_candidates = [max_image_dpi, max(max_image_dpi - 50, 50)]
                else:
                    if level_key == "low":
                        dpi_candidates = [300, 200, 150]
                    elif level_key == "medium":
                        dpi_candidates = [150, 96, 72]
                    else:
                        dpi_candidates = [96, 72, 50]

                satisfied = False
                for setting in gs_settings_order.get(level_key, ["/ebook"]):
                    for dpi in dpi_candidates:
                        size_after = _run_gs(setting, dpi)
                        if size_after:
                            used_method = "ghostscript"
                            if desired_size is None or size_after <= desired_size:
                                satisfied = True
                                break
                    if satisfied:
                        break

            if not used_method:
                levels = {
                    "low": {"garbage": 4},
                    "medium": {"garbage": 3},
                    "high": {"garbage": 2},
                }
                level = levels.get(level_key, levels["medium"])
                doc = fitz.open(input_path)
                doc.save(output_path, deflate=True, garbage=level["garbage"], clean=True)
                doc.close()
                used_method = "pymupdf"

            try:
                if original_size is not None and os.path.exists(output_path):
                    new_size = os.path.getsize(output_path)
                    if desired_size and new_size > desired_size and gs_exe:
                        size_after = _run_gs("/screen", 72)
                        if size_after and size_after <= desired_size:
                            used_method = "ghostscript"
                    elif new_size >= original_size and used_method != "pymupdf":
                        doc = fitz.open(input_path)
                        doc.save(output_path, deflate=True, garbage=2, clean=True)
                        doc.close()
                        used_method = "pymupdf"
            except Exception:
                pass

            return output_path
        except Exception as e:
            raise FileProcessingError(f"Error compressing PDF: {str(e)}")

    @staticmethod
    def split_pdf(
        input_path: str,
        split_type: str = "every_page",
        ranges: Optional[List[str]] = None,
        output_prefix: Optional[str] = None,
        zip_output: bool = False,
    ) -> Dict[str, Any]:
        try:
            base_name = os.path.splitext(os.path.basename(input_path))[0] or "pdf"
            prefix_base = (output_prefix or base_name)
            sanitized_prefix = re.sub(r"[^A-Za-z0-9._-]+", "_", prefix_base).strip("._") or "pdf"

            output_root = settings.output_dir if settings and settings.output_dir else os.path.dirname(input_path)
            os.makedirs(output_root, exist_ok=True)
            folder_name = sanitized_prefix
            output_folder = os.path.join(output_root, folder_name)
            os.makedirs(output_folder, exist_ok=True)

            reader = PdfReader(input_path)
            total_pages = len(reader.pages)
            results: List[Dict[str, Any]] = []

            def _unique_path(name_without_ext: str) -> str:
                candidate = f"{name_without_ext}.pdf"
                path = os.path.join(output_folder, candidate)
                counter = 1
                while os.path.exists(path):
                    candidate = f"{name_without_ext}_{counter}.pdf"
                    path = os.path.join(output_folder, candidate)
                    counter += 1
                return path

            def _emit(start: int, end: int, pages_desc: str, pages_list: List[int]):
                writer = PdfWriter()
                for p in pages_list:
                    if p < 1 or p > total_pages:
                        raise FileProcessingError("Invalid page range")
                    writer.add_page(reader.pages[p - 1])
                name_without_ext = f"{sanitized_prefix}_{pages_desc}" if pages_desc else f"{sanitized_prefix}"
                out_path = _unique_path(name_without_ext)
                with open(out_path, "wb") as f:
                    writer.write(f)
                results.append({
                    "path": out_path,
                    "filename": os.path.basename(out_path),
                    "pages": pages_list,
                })

            st = (split_type or "").strip().lower()
            if st == "every_page" or (st == "" and not ranges):
                for i in range(1, total_pages + 1):
                    _emit(i, i, f"page_{i}", [i])
            elif st == "page_ranges" or (st == "" and ranges):
                if not ranges:
                    raise FileProcessingError("page_ranges required for split_type=page_ranges")
                for token in ranges:
                    token = (token or "").strip()
                    if not token:
                        continue
                    if "-" in token:
                        s, e = token.split("-", 1)
                        start = int(s)
                        end = int(e)
                    else:
                        start = int(token)
                        end = start
                    if start > end:
                        start, end = end, start
                    pages_list = list(range(start, end + 1))
                    # Naming as requested: prefix_page_1 and prefix_page_3_5
                    desc = f"page_{start}_{end}" if start != end else f"page_{start}"
                    _emit(start, end, desc, pages_list)
            else:
                raise FileProcessingError("Unsupported split type")

            result: Dict[str, Any] = {"files": results, "count": len(results), "folder_name": folder_name}
            if zip_output:
                zip_name = f"{folder_name}.zip"
                zip_path = os.path.join(output_root, zip_name)
                counter = 1
                while os.path.exists(zip_path):
                    zip_name = f"{folder_name}_{counter}.zip"
                    zip_path = os.path.join(output_root, zip_name)
                    counter += 1
                with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                    for item in results:
                        zf.write(item["path"], arcname=os.path.join(folder_name, item["filename"]))
                result["zip_path"] = zip_path
                result["zip_filename"] = os.path.basename(zip_path)
            return result
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def extract_pages_to_single(input_path: str, output_path: str, ranges: List[str]) -> str:
        try:
            reader = PdfReader(input_path)
            writer = PdfWriter()
            total_pages = len(reader.pages)

            for token in ranges:
                token = (token or "").strip()
                if not token:
                    continue
                if "-" in token:
                    s, e = token.split("-", 1)
                    start = int(s)
                    end = int(e)
                else:
                    start = int(token)
                    end = start
                if start < 1 or end < 1 or start > total_pages or end > total_pages:
                    raise FileProcessingError("Invalid page range")
                if start > end:
                    start, end = end, start
                for p in range(start, end + 1):
                    writer.add_page(reader.pages[p - 1])

            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path, "wb") as f:
                writer.write(f)
            return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def remove_pages(input_path: str, output_path: str, pages_to_remove: List[int]) -> str:
        try:
            reader = PdfReader(input_path)
            total_pages = len(reader.pages)
            pages_set = set(int(p) for p in pages_to_remove if p is not None)
            for p in pages_set:
                if p < 1 or p > total_pages:
                    raise FileProcessingError("Invalid page number")
            keep = [i for i in range(1, total_pages + 1) if i not in pages_set]
            if not keep:
                raise FileProcessingError("No pages left after removal")
            writer = PdfWriter()
            for i in keep:
                writer.add_page(reader.pages[i - 1])
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path, "wb") as f:
                writer.write(f)
            return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def extract_pages(input_path: str, output_path: str, pages_to_extract: List[int]) -> str:
        try:
            reader = PdfReader(input_path)
            total_pages = len(reader.pages)
            ordered = []
            seen = set()
            for p in pages_to_extract:
                p = int(p)
                if p in seen:
                    continue
                seen.add(p)
                if p < 1 or p > total_pages:
                    raise FileProcessingError("Invalid page number")
                ordered.append(p)
            if not ordered:
                raise FileProcessingError("No pages selected")
            writer = PdfWriter()
            for p in ordered:
                writer.add_page(reader.pages[p - 1])
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path, "wb") as f:
                writer.write(f)
            return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def rotate_pdf(input_path: str, output_path: str, rotation: int = 90) -> str:
        try:
            deg = int(rotation) if rotation is not None else 90
            deg = deg % 360
            if deg % 90 != 0:
                deg = (deg // 90) * 90

            try:
                doc = fitz.open(input_path)
                if deg == 0:
                    os.makedirs(os.path.dirname(output_path), exist_ok=True)
                    doc.save(output_path)
                    doc.close()
                    return output_path

                for page in doc:
                    try:
                        if hasattr(page, "set_rotation"):
                            page.set_rotation(deg)
                        elif hasattr(page, "setRotation"):
                            page.setRotation(deg)
                        else:
                            current = 0
                            try:
                                current = page.rotation  # type: ignore
                            except Exception:
                                current = 0
                            target = (current + deg) % 360
                            if hasattr(page, "set_rotation"):
                                page.set_rotation(target)
                    except Exception:
                        pass

                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                doc.save(output_path)
                doc.close()
                return output_path
            except Exception:
                reader = PdfReader(input_path)
                writer = PdfWriter()
                if deg == 0:
                    for page in reader.pages:
                        writer.add_page(page)
                else:
                    for page in reader.pages:
                        try:
                            if hasattr(page, "rotate_clockwise"):
                                rotated = page.rotate_clockwise(deg)
                            elif hasattr(page, "rotateCounterClockwise") and deg < 0:
                                rotated = page.rotateCounterClockwise(abs(deg))
                            else:
                                rotated = page.rotateClockwise(deg)
                        except Exception:
                            rotated = page
                        writer.add_page(rotated)
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, "wb") as f:
                    writer.write(f)
                return output_path
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def add_watermark(
        input_path: str,
        output_path: str,
        watermark_text: str,
        position: str = "center",
    ) -> str:
        try:
            if not watermark_text or not str(watermark_text).strip():
                raise FileProcessingError("Watermark text cannot be empty")

            pos = (position or "center").strip().lower()
            base_positions = {
                "top-left", "top-center", "top-right",
                "middle-left", "center", "middle-right",
                "bottom-left", "bottom-center", "bottom-right",
            }
            diag = False
            base = pos
            if pos.endswith("-diagonal"):
                diag = True
                base = pos[:-9]
            if base not in base_positions:
                base = "center"
                diag = True

            # Primary implementation: draw watermark using PyMuPDF
            try:
                doc = fitz.open(input_path)

                def _rect_for_position(r: fitz.Rect, where: str) -> fitz.Rect:
                    w, h = r.width, r.height
                    margin = max(20, int(min(w, h) * 0.03))
                    box_w = max(200, int(w * 0.6))
                    box_h = max(60, int(h * 0.15))
                    if where == "top-left":
                        return fitz.Rect(r.x0 + margin, r.y0 + margin, r.x0 + margin + box_w, r.y0 + margin + box_h)
                    if where == "top-center":
                        cx = r.x0 + w / 2
                        return fitz.Rect(cx - box_w / 2, r.y0 + margin, cx + box_w / 2, r.y0 + margin + box_h)
                    if where == "top-right":
                        return fitz.Rect(r.x1 - margin - box_w, r.y0 + margin, r.x1 - margin, r.y0 + margin + box_h)
                    if where == "middle-left":
                        cy = r.y0 + h / 2
                        return fitz.Rect(r.x0 + margin, cy - box_h / 2, r.x0 + margin + box_w, cy + box_h / 2)
                    if where == "middle-right":
                        cy = r.y0 + h / 2
                        return fitz.Rect(r.x1 - margin - box_w, cy - box_h / 2, r.x1 - margin, cy + box_h / 2)
                    if where == "bottom-left":
                        return fitz.Rect(r.x0 + margin, r.y1 - margin - box_h, r.x0 + margin + box_w, r.y1 - margin)
                    if where == "bottom-center":
                        cx = r.x0 + w / 2
                        return fitz.Rect(cx - box_w / 2, r.y1 - margin - box_h, cx + box_w / 2, r.y1 - margin)
                    if where == "bottom-right":
                        return fitz.Rect(r.x1 - margin - box_w, r.y1 - margin - box_h, r.x1 - margin, r.y1 - margin)
                    # center
                    cx = r.x0 + w / 2
                    cy = r.y0 + h / 2
                    return fitz.Rect(cx - box_w / 2, cy - box_h / 2, cx + box_w / 2, cy + box_h / 2)

                for page in doc:
                    r = page.rect
                    target_rect = _rect_for_position(r, base)
                    fs = max(24, int(r.width * 0.06))
                    try:
                        page.insert_textbox(
                            target_rect,
                            str(watermark_text),
                            fontsize=fs,
                            fontname="Helvetica-Bold",
                            color=(0.5, 0.5, 0.5),
                            align=fitz.TEXT_ALIGN_CENTER,
                            rotate=(45 if diag or base == "center" else 0),
                        )
                    except Exception:
                        # Fallback text insert without textbox
                        center_pt = fitz.Point(target_rect.x0 + target_rect.width / 2, target_rect.y0 + target_rect.height / 2)
                        page.insert_text(
                            center_pt,
                            str(watermark_text),
                            fontsize=fs,
                            fontname="Helvetica-Bold",
                            color=(0.5, 0.5, 0.5),
                        )

                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                doc.save(output_path)
                doc.close()
                return output_path
            except Exception:
                # Fallback: overlay watermark using ReportLab + PyPDF2
                reader = PdfReader(input_path)
                writer = PdfWriter()

                for page in reader.pages:
                    width = float(page.mediabox.width)
                    height = float(page.mediabox.height)

                    buf = BytesIO()
                    c = canvas.Canvas(buf, pagesize=(width, height))
                    try:
                        c.setFillColorRGB(0.5, 0.5, 0.5)
                        try:
                            # Transparency if available
                            c.setFillAlpha(0.2)
                        except Exception:
                            pass
                        c.setFont("Helvetica-Bold", max(24, int(width * 0.06)))
                        margin = max(20, int(min(width, height) * 0.03))
                        box_w = max(200, int(width * 0.6))
                        box_h = max(60, int(height * 0.15))
                        def _rect_for_position_fallback(where: str):
                            if where == "top-left":
                                return (margin + box_w / 2, margin + box_h / 2)
                            if where == "top-center":
                                return (width / 2.0, margin + box_h / 2)
                            if where == "top-right":
                                return (width - margin - box_w / 2, margin + box_h / 2)
                            if where == "middle-left":
                                return (margin + box_w / 2, height / 2.0)
                            if where == "middle-right":
                                return (width - margin - box_w / 2, height / 2.0)
                            if where == "bottom-left":
                                return (margin + box_w / 2, height - margin - box_h / 2)
                            if where == "bottom-center":
                                return (width / 2.0, height - margin - box_h / 2)
                            if where == "bottom-right":
                                return (width - margin - box_w / 2, height - margin - box_h / 2)
                            return (width / 2.0, height / 2.0)

                        cx, cy = _rect_for_position_fallback(base)
                        if diag or base == "center":
                            c.saveState()
                            c.translate(cx, cy)
                            c.rotate(45)
                            c.drawCentredString(0, 0, str(watermark_text))
                            c.restoreState()
                        else:
                            if base.endswith("left"):
                                c.drawString(cx - box_w / 2, cy, str(watermark_text))
                            elif base.endswith("right"):
                                c.drawRightString(cx + box_w / 2, cy, str(watermark_text))
                            else:
                                c.drawCentredString(cx, cy, str(watermark_text))
                    finally:
                        c.save()
                    buf.seek(0)
                    wm_pdf = PdfReader(buf)
                    try:
                        page.merge_page(wm_pdf.pages[0])
                    except Exception:
                        # Some PyPDF2 versions use mergePage
                        try:
                            page.mergePage(wm_pdf.pages[0])
                        except Exception:
                            pass
                    writer.add_page(page)

                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, "wb") as f:
                    writer.write(f)
                return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))

    @staticmethod
    def add_page_numbers(
        input_path: str,
        output_path: str,
        position: str = "bottom-center",
        start_page: int = 1,
        fmt: str = "{page}",
        font_size: float = 12.0,
    ) -> str:
        try:
            pos = (position or "bottom-center").strip().lower()
            allowed = {
                "top-left", "top-center", "top-right",
                "bottom-left", "bottom-center", "bottom-right",
            }
            if pos not in allowed:
                pos = "bottom-center"

            sp = int(start_page) if start_page is not None else 1
            if sp < 1:
                sp = 1
            fs = float(font_size) if font_size is not None else 12.0
            if fs <= 0:
                fs = 12.0

            # Primary implementation: PyMuPDF
            try:
                doc = fitz.open(input_path)
                for idx, page in enumerate(doc, start=1):
                    if idx < sp:
                        continue
                    r = page.rect
                    margin = max(20, int(min(r.width, r.height) * 0.03))
                    box_h = max(24, int(fs * 2))
                    # Rectangle spanning width (respect margins), choose top/bottom band
                    if pos.startswith("top-"):
                        rect = fitz.Rect(r.x0 + margin, r.y0 + margin, r.x1 - margin, r.y0 + margin + box_h)
                    else:
                        rect = fitz.Rect(r.x0 + margin, r.y1 - margin - box_h, r.x1 - margin, r.y1 - margin)

                    align_map = {
                        "left": fitz.TEXT_ALIGN_LEFT,
                        "center": fitz.TEXT_ALIGN_CENTER,
                        "right": fitz.TEXT_ALIGN_RIGHT,
                    }
                    align_key = pos.split("-")[1]  # left/center/right
                    align = align_map.get(align_key, fitz.TEXT_ALIGN_CENTER)

                    text = (fmt or "{page}")
                    try:
                        text = text.replace("{page}", str(idx))
                    except Exception:
                        text = str(idx)

                    page.insert_textbox(
                        rect,
                        text,
                        fontsize=fs,
                        fontname="Helvetica-Bold",
                        color=(0, 0, 0),
                        align=align,
                    )

                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                doc.save(output_path)
                doc.close()
                return output_path
            except Exception:
                # Fallback: ReportLab overlay merged via PyPDF2
                reader = PdfReader(input_path)
                writer = PdfWriter()
                for idx, page in enumerate(reader.pages, start=1):
                    width = float(page.mediabox.width)
                    height = float(page.mediabox.height)
                    buf = BytesIO()
                    c = canvas.Canvas(buf, pagesize=(width, height))
                    try:
                        c.setFont("Helvetica-Bold", fs)
                        margin = max(20, int(min(width, height) * 0.03))
                        text = (fmt or "{page}").replace("{page}", str(idx))
                        # Compute anchor x,y
                        if pos.startswith("top-"):
                            y = margin + fs
                        else:
                            y = height - margin - fs
                        key = pos.split("-")[1]
                        if key == "left":
                            x = margin
                            c.drawString(x, y, text)
                        elif key == "right":
                            x = width - margin
                            c.drawRightString(x, y, text)
                        else:
                            x = width / 2.0
                            c.drawCentredString(x, y, text)
                    finally:
                        c.save()
                    buf.seek(0)
                    wm_pdf = PdfReader(buf)
                    try:
                        page.merge_page(wm_pdf.pages[0])
                    except Exception:
                        try:
                            page.mergePage(wm_pdf.pages[0])
                        except Exception:
                            pass
                    writer.add_page(page)

                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, "wb") as f:
                    writer.write(f)
                return output_path
        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(str(e))
