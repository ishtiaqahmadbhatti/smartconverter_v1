"""
Simplified Website Conversion Service

This service handles website and HTML conversion operations with minimal dependencies.
"""

import os
import json
import tempfile
import logging
from typing import Dict, Any, Optional
from pathlib import Path
import base64
import io
import uuid

# Basic HTML/Web conversion libraries
from markdown import markdown
from bs4 import BeautifulSoup
import requests
from PIL import Image, ImageDraw, ImageFont

# Selenium for professional rendering
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service as ChromeService
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.print_page_options import PrintOptions

# Document processing libraries
import docx
from pptx import Presentation
import pandas as pd
import fitz  # PyMuPDF

# Database logging
from app.services.request_logging_service import RequestLoggingService

logger = logging.getLogger(__name__)


class WebsiteConversionService:
    """Service for website and HTML conversion operations."""
    
    @staticmethod
    def _get_chrome_options():
        """Get Chrome options for headless execution."""
        chrome_options = Options()
        chrome_options.add_argument('--headless')
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument('--disable-gpu')
        chrome_options.add_argument('--disable-gpu')
        return chrome_options

    @staticmethod
    def _generate_error_pdf(output_path: str, error_message: str, url_or_content: str):
        """Generate a PDF with error details using PIL."""
        try:
            # Create an image with the error message
            width, height = 1240, 1754  # A4 size at ~150 DPI
            img = Image.new('RGB', (width, height), color=(255, 255, 255))
            d = ImageDraw.Draw(img)
            
            # Try to load a font, fallback to default
            try:
                # Big font for header
                font_header = ImageFont.truetype("arial.ttf", 40)
                # Medium font for content
                font_text = ImageFont.truetype("arial.ttf", 20)
            except IOError:
                font_header = ImageFont.load_default()
                font_text = ImageFont.load_default()

            # Draw Header
            d.text((50, 50), "PDF Conversion Failed", fill=(255, 0, 0), font=font_header)
            
            # Draw Source Info
            source_text = f"Source: {url_or_content[:100]}..." if len(url_or_content) > 100 else f"Source: {url_or_content}"
            d.text((50, 120), source_text, fill=(0, 0, 0), font=font_text)
            
            # Draw Error Message
            import textwrap
            wrapped_error = textwrap.fill(f"Error Details: {error_message}", width=80)
            d.text((50, 160), wrapped_error, fill=(255, 0, 0), font=font_text)
            
            # Draw Note
            note = "Note: Please ensure Google Chrome is correctly installed and configured on the server."
            d.text((50, 400), note, fill=(100, 100, 100), font=font_text)
            
            # Draw Timestamp
            timestamp = f"Generated: {uuid.uuid4()}"
            d.text((50, 450), timestamp, fill=(100, 100, 100), font=font_text)

            # Save as PDF
            img.save(output_path, "PDF", resolution=100.0)
            return output_path
            
        except Exception as e:
            logger.error(f"Failed to generate error PDF: {str(e)}")
            # Last resort: create an empty text file or simple PDF if possible, but for now just re-raise
            raise e

    @staticmethod
    def convert_html_file_to_pdf(file_path: str, output_filename: str = None) -> str:
        """Convert HTML file to PDF using Selenium (Professional)."""
        driver = None
        try:
            # Create unique filename if not provided
            if output_filename:
                if not output_filename.lower().endswith('.pdf'):
                    output_filename += '.pdf'
                filename = output_filename
            else:
                unique_id = str(uuid.uuid4())
                filename = f"html_to_pdf_{unique_id}.pdf"
            
            output_path = os.path.join("outputs", filename)
            os.makedirs("outputs", exist_ok=True)

            # Initialize WebDriver
            driver = webdriver.Chrome(service=ChromeService(ChromeDriverManager().install()), options=WebsiteConversionService._get_chrome_options())
            
            # Load HTML file
            abs_path = os.path.abspath(file_path)
            driver.get(f"file:///{abs_path}")
            
            # Print to PDF
            print_options = PrintOptions()
            print_options.background = True
            pdf_data = driver.print_page(print_options)
            
            # Save PDF
            with open(output_path, 'wb') as f:
                f.write(base64.b64decode(pdf_data))
                
            return output_path

        except Exception as e:
            logger.error(f"Error converting HTML file to PDF: {str(e)}")
            try:
                # Generate fallback PDF
                return WebsiteConversionService._generate_error_pdf(output_path, str(e), f"File: {os.path.basename(file_path)}")
            except Exception as fallback_error:
                raise Exception(f"Failed to convert HTML file to PDF and fallback failed: {str(e)}")
        finally:
            if driver:
                driver.quit()

    @staticmethod
    def html_to_pdf(html_content: str, css_content: str = None, output_filename: str = None) -> str:
        """Convert HTML content to PDF using Selenium (Professional)."""
        driver = None
        html_file_path = None
        try:
            # Create unique filename if not provided
            if output_filename:
                if not output_filename.lower().endswith('.pdf'):
                    output_filename += '.pdf'
                filename = output_filename
            else:
                unique_id = str(uuid.uuid4())
                filename = f"html_to_pdf_{unique_id}.pdf"
            
            output_path = os.path.join("outputs", filename)
            os.makedirs("outputs", exist_ok=True)

            # Create temporary HTML file
            # Inject CSS if provided
            if css_content:
                html_content = f"<style>{css_content}</style>\n{html_content}"
                
            with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as html_file:
                html_file.write(html_content)
                html_file_path = html_file.name

            # Initialize WebDriver
            driver = webdriver.Chrome(service=ChromeService(ChromeDriverManager().install()), options=WebsiteConversionService._get_chrome_options())
            
            # Load HTML file
            # Use absolute path with file protocol
            abs_path = os.path.abspath(html_file_path)
            driver.get(f"file:///{abs_path}")
            
            # Print to PDF
            print_options = PrintOptions()
            print_options.background = True
            pdf_data = driver.print_page(print_options)
            
            # Save PDF
            with open(output_path, 'wb') as f:
                f.write(base64.b64decode(pdf_data))
                
            return output_path

        except Exception as e:
            logger.error(f"Error converting HTML to PDF: {str(e)}")
            try:
                # Generate fallback PDF
                preview_content = html_content[:50] + "..." if html_content else "Empty Content"
                return WebsiteConversionService._generate_error_pdf(output_path, str(e), f"HTML Content: {preview_content}")
            except Exception as fallback_error:
                raise Exception(f"Failed to convert HTML to PDF and fallback failed: {str(e)}")
        finally:
            if driver:
                driver.quit()
            if html_file_path and os.path.exists(html_file_path):
                try:
                    os.unlink(html_file_path)
                except:
                    pass

    @staticmethod
    def website_to_pdf(url: str, output_filename: str = None) -> str:
        """Convert Website URL to PDF using Selenium (Professional)."""
        driver = None
        try:
            # Create unique filename if not provided
            if output_filename:
                if not output_filename.lower().endswith('.pdf'):
                    output_filename += '.pdf'
                filename = output_filename
            else:
                unique_id = str(uuid.uuid4())
                filename = f"website_to_pdf_{unique_id}.pdf"
            
            output_path = os.path.join("outputs", filename)
            os.makedirs("outputs", exist_ok=True)

            # Initialize WebDriver
            driver = webdriver.Chrome(service=ChromeService(ChromeDriverManager().install()), options=WebsiteConversionService._get_chrome_options())
            
            # Load URL
            driver.get(url)
            
            # Print to PDF
            print_options = PrintOptions()
            print_options.background = True
            pdf_data = driver.print_page(print_options)
            
            # Save PDF
            with open(output_path, 'wb') as f:
                f.write(base64.b64decode(pdf_data))
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Website to PDF: {str(e)}")
            try:
                # Generate fallback PDF
                return WebsiteConversionService._generate_error_pdf(output_path, str(e), url)
            except Exception as fallback_error:
                raise Exception(f"Failed to convert Website to PDF and fallback failed: {str(e)}")
        finally:
            if driver:
                driver.quit()

    @staticmethod
    def word_to_html(file_content: bytes, original_filename: str, output_filename: str = None) -> str:
        """Convert Word document to HTML with comprehensive formatting."""
        try:
            # Determine title
            title = os.path.splitext(original_filename)[0] if original_filename else "Converted Document"
            if output_filename and output_filename.strip() and output_filename.lower() != "string":
                title = os.path.splitext(output_filename)[0]

            # Determine filename
            if output_filename and output_filename.strip() and output_filename.lower() != "string":
                if not output_filename.lower().endswith('.html'):
                    output_filename += '.html'
                filename = output_filename
            else:
                if original_filename:
                    base_name = os.path.splitext(original_filename)[0]
                    filename = f"{base_name}.html"
                else:
                    unique_id = str(uuid.uuid4())
                    filename = f"word_to_html_{unique_id}.html"
            
            output_path = os.path.join("outputs", filename)
            os.makedirs("outputs", exist_ok=True)

            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as word_file:
                word_file.write(file_content)
                word_file_path = word_file.name
            
            # Read Word document
            doc = docx.Document(word_file_path)
            
            # CSS for better styling
            css = """
            <style>
                body { font-family: Calibri, Arial, sans-serif; line-height: 1.5; padding: 20px; max-width: 800px; margin: 0 auto; background-color: #fff; color: #000; }
                table { border-collapse: collapse; width: 100%; margin-bottom: 20px; }
                td, th { border: 1px solid #a0a0a0; padding: 8px; vertical-align: top; }
                tr:nth-child(even) { background-color: #f9f9f9; }
                th { background-color: #f2f2f2; font-weight: bold; text-align: left; }
                p { margin-bottom: 10px; margin-top: 0; }
                ul, ol { margin-bottom: 10px; padding-left: 40px; }
                li { margin-bottom: 5px; }
                a { color: #0563c1; text-decoration: underline; }
                .highlight { background-color: yellow; }
            </style>
            """
            
            html_content = f"<html><head><title>{title}</title>{css}</head><body>"
            
            # Helper to process runs
            def process_run(run):
                styles = []
                if run.bold: styles.append("font-weight: bold")
                if run.italic: styles.append("font-style: italic")
                if run.underline: styles.append("text-decoration: underline")
                if run.font.strike: styles.append("text-decoration: line-through")
                if run.font.subscript: styles.append("vertical-align: sub; font-size: smaller")
                if run.font.superscript: styles.append("vertical-align: super; font-size: smaller")
                
                if run.font.color and run.font.color.rgb:
                    styles.append(f"color: #{run.font.color.rgb}")
                
                if run.font.highlight_color:
                    styles.append("background-color: yellow") # Simplified highlight
                
                if run.font.size:
                    size = run.font.size.pt
                    styles.append(f"font-size: {size}pt")
                
                text = run.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                
                # Handle images (basic support)
                if 'graphic' in run._element.xml:
                    # This is complex in python-docx, skipping for stability unless requested specifically
                    # Ideally we would extract blip relationships
                    pass

                if styles:
                    return f"<span style='{'; '.join(styles)}'>{text}</span>"
                else:
                    return text

            # Helper to process paragraph content (including hyperlinks)
            def get_paragraph_content(paragraph):
                content = ""
                # We need to iterate over the XML children to handle hyperlinks correctly
                from docx.oxml.text.run import CT_R
                from docx.oxml.text.hyperlink import CT_Hyperlink
                
                for child in paragraph._element:
                    if isinstance(child, CT_R):
                        # It's a Run
                        from docx.text.run import Run
                        run = Run(child, paragraph)
                        content += process_run(run)
                    elif isinstance(child, CT_Hyperlink):
                        # It's a Hyperlink
                        # Get the relationship ID
                        try:
                            rId = child.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                            if rId:
                                url = paragraph.part.rels[rId].target_ref
                                link_text = ""
                                for sub_child in child:
                                    if isinstance(sub_child, CT_R):
                                        run = Run(sub_child, paragraph)
                                        link_text += process_run(run)
                                content += f"<a href='{url}'>{link_text}</a>"
                        except:
                            pass # Skip broken links
                return content

            # Helper to process paragraph
            def process_paragraph(paragraph):
                align_map = { 0: 'left', 1: 'center', 2: 'right', 3: 'justify' }
                align = align_map.get(paragraph.alignment, 'left')
                style = f"text-align: {align};"
                
                content = get_paragraph_content(paragraph)
                
                if not content.strip():
                    return ""

                # Handle Lists
                # This is a simple heuristic. For robust lists, we need to track numbering state.
                # Here we just check if it has numbering properties.
                is_list = False
                if paragraph._element.pPr is not None and paragraph._element.pPr.numPr is not None:
                    is_list = True
                
                tag = "p"
                if paragraph.style.name.startswith('Heading'):
                    level = paragraph.style.name.split()[-1] if paragraph.style.name.split()[-1].isdigit() else '1'
                    tag = f"h{level}"
                
                if is_list:
                    # For now, render as list item, but we need wrapping ul/ol.
                    # Since we are processing linearly, we can't easily wrap. 
                    # We'll use a div with display: list-item or just a bullet char for simplicity in this pass.
                    # Or better: return a special marker tuple?
                    # Let's stick to simple HTML: use <li> but we might miss the <ul> wrapper.
                    # Browsers often handle orphan <li> but it's invalid.
                    # Let's use a styled div to simulate list item if we can't wrap.
                    return f"<div style='display: list-item; margin-left: 20px; {style}'>{content}</div>"
                
                return f"<{tag} style='{style}'>{content}</{tag}>"

            # Main processing loop
            for child in doc.element.body:
                if child.tag.endswith('p'):
                    from docx.text.paragraph import Paragraph
                    para = Paragraph(child, doc)
                    html_content += process_paragraph(para)
                elif child.tag.endswith('tbl'):
                    from docx.table import Table
                    table = Table(child, doc)
                    html_content += "<table>"
                    for row in table.rows:
                        html_content += "<tr>"
                        for cell in row.cells:
                            # Handle colspan (grid_span)
                            # python-docx doesn't expose grid_span easily on _Cell, need xml
                            tcPr = cell._element.tcPr
                            colspan = 1
                            if tcPr is not None:
                                grid_span = tcPr.find(f"{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}gridSpan")
                                if grid_span is not None:
                                    try:
                                        val = grid_span.get(f"{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}val")
                                        # Handle potential float strings or invalid values
                                        colspan = int(float(val)) if val else 1
                                    except Exception:
                                        colspan = 1
                            
                            attr = f" colspan='{colspan}'" if colspan > 1 else ""
                            
                            html_content += f"<td{attr}>"
                            for para in cell.paragraphs:
                                html_content += process_paragraph(para)
                            html_content += "</td>"
                        html_content += "</tr>"
                    html_content += "</table>"

            html_content += "</body></html>"
            
            # Save HTML to file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([word_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Word to HTML: {str(e)}")
            raise Exception(f"Failed to convert Word to HTML: {str(e)}")
    
    @staticmethod
    def powerpoint_to_html(file_content: bytes, original_filename: str, output_filename: str = None) -> str:
        """Convert PowerPoint presentation to HTML."""
        try:
            # Determine title
            title = os.path.splitext(original_filename)[0] if original_filename else "Converted Presentation"
            if output_filename and output_filename.strip() and output_filename.lower() != "string":
                title = os.path.splitext(output_filename)[0]

            # Determine filename
            if output_filename and output_filename.strip() and output_filename.lower() != "string":
                if not output_filename.lower().endswith('.html'):
                    output_filename += '.html'
                filename = output_filename
            else:
                if original_filename:
                    base_name = os.path.splitext(original_filename)[0]
                    filename = f"{base_name}.html"
                else:
                    unique_id = str(uuid.uuid4())
                    filename = f"ppt_to_html_{unique_id}.html"
            
            output_path = os.path.join("outputs", filename)
            os.makedirs("outputs", exist_ok=True)

            # Create temporary file for PowerPoint
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as ppt_file:
                ppt_file.write(file_content)
                ppt_file_path = ppt_file.name
            
            # Read PowerPoint presentation
            prs = Presentation(ppt_file_path)
            
            # CSS for better styling
            css = """
            <style>
                body { font-family: Arial, sans-serif; line-height: 1.6; padding: 20px; max-width: 1000px; margin: 0 auto; background-color: #f0f0f0; color: #333; }
                .slide { background-color: white; padding: 40px; margin-bottom: 40px; box-shadow: 0 4px 8px rgba(0,0,0,0.1); border-radius: 8px; aspect-ratio: 16/9; display: flex; flex-direction: column; overflow: hidden; position: relative; }
                .slide-number { position: absolute; bottom: 20px; right: 20px; color: #888; font-size: 12px; }
                h1, h2, h3 { color: #2c3e50; margin-top: 0; }
                ul { padding-left: 20px; }
                li { margin-bottom: 8px; }
                p { margin-bottom: 10px; }
                .notes { background-color: #fff9c4; padding: 10px; border-left: 4px solid #fbc02d; margin-top: 20px; font-size: 0.9em; }
            </style>
            """
            
            html_content = f"<html><head><title>{title}</title>{css}</head><body>"
            
            for i, slide in enumerate(prs.slides, 1):
                html_content += f"<div class='slide'><div class='slide-number'>Slide {i}</div>"
                
                # Extract title
                if slide.shapes.title:
                    html_content += f"<h2>{slide.shapes.title.text}</h2>"
                
                # Extract text content
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    # Skip title as we already handled it
                    if shape == slide.shapes.title:
                        continue
                        
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            # Basic bullet detection
                            if paragraph.level > 0:
                                html_content += f"<ul style='margin-left: {paragraph.level * 20}px'><li>{text}</li></ul>"
                            else:
                                html_content += f"<p>{text}</p>"
                
                html_content += "</div>"
                
                # Extract notes if any
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        html_content += f"<div class='notes'><strong>Notes:</strong><br>{notes}</div>"
                
                html_content += "<hr style='border: 0; clear: both; margin-bottom: 40px;'>"
            
            html_content += "</body></html>"
            
            # Save HTML to file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([ppt_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting PowerPoint to HTML: {str(e)}")
            raise Exception(f"Failed to convert PowerPoint to HTML: {str(e)}")
    
    @staticmethod
    def markdown_to_html(markdown_content: str, original_filename: str = None, output_filename: str = None) -> str:
        """Convert Markdown content to HTML."""
        try:
            # Determine title
            title = os.path.splitext(original_filename)[0] if original_filename else "Converted Markdown"
            if output_filename and output_filename.strip() and output_filename.lower() != "string":
                title = os.path.splitext(output_filename)[0]

            # Determine filename
            if output_filename and output_filename.strip() and output_filename.lower() != "string":
                if not output_filename.lower().endswith('.html'):
                    output_filename += '.html'
                filename = output_filename
            else:
                if original_filename:
                    base_name = os.path.splitext(original_filename)[0]
                    filename = f"{base_name}.html"
                else:
                    unique_id = str(uuid.uuid4())
                    filename = f"markdown_to_html_{unique_id}.html"
            
            output_path = os.path.join("outputs", filename)
            os.makedirs("outputs", exist_ok=True)

            html_content = markdown(markdown_content, extensions=['tables', 'fenced_code', 'codehilite'])
            
            # Wrap in proper HTML structure
            full_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <title>{title}</title>
                <style>
                    body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 0 auto; padding: 40px; max-width: 900px; color: #333; background-color: #fff; }}
                    h1, h2, h3, h4, h5, h6 {{ color: #2c3e50; margin-top: 24px; margin-bottom: 16px; }}
                    h1 {{ border-bottom: 1px solid #eee; padding-bottom: 0.3em; }}
                    h2 {{ border-bottom: 1px solid #eee; padding-bottom: 0.3em; }}
                    p {{ margin-bottom: 16px; }}
                    a {{ color: #0366d6; text-decoration: none; }}
                    a:hover {{ text-decoration: underline; }}
                    code {{ background-color: #f6f8fa; padding: 0.2em 0.4em; border-radius: 3px; font-family: SFMono-Regular, Consolas, "Liberation Mono", Menlo, monospace; font-size: 85%; }}
                    pre {{ background-color: #f6f8fa; padding: 16px; border-radius: 6px; overflow: auto; line-height: 1.45; }}
                    pre code {{ background-color: transparent; padding: 0; font-size: 100%; }}
                    blockquote {{ border-left: 4px solid #dfe2e5; padding: 0 1em; color: #6a737d; margin: 0 0 16px 0; }}
                    table {{ border-collapse: collapse; width: 100%; margin-bottom: 16px; }}
                    th, td {{ border: 1px solid #dfe2e5; padding: 6px 13px; }}
                    th {{ background-color: #f6f8fa; font-weight: 600; }}
                    tr:nth-child(2n) {{ background-color: #f8f8f8; }}
                    img {{ max-width: 100%; box-sizing: content-box; background-color: #fff; }}
                    hr {{ height: 0.25em; padding: 0; margin: 24px 0; background-color: #e1e4e8; border: 0; }}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
            
            # Save HTML to file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_html)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Markdown to HTML: {str(e)}")
            raise Exception(f"Failed to convert Markdown to HTML: {str(e)}")
    
    @staticmethod
    def html_table_to_csv(html_content: str, output_filename: str = None, original_filename: str = None) -> str:
        """Convert HTML table to CSV."""
        # Force reload
        try:
            # Determine filename
            if output_filename and output_filename.strip() and output_filename.lower() != "string":
                if not output_filename.lower().endswith('.csv'):
                    output_filename += '.csv'
                filename = output_filename
            elif original_filename:
                base_name = os.path.splitext(original_filename)[0]
                filename = f"{base_name}.csv"
            else:
                unique_id = str(uuid.uuid4())
                filename = f"html_table_to_csv_{unique_id}.csv"
            
            output_path = os.path.join("outputs", filename)
            os.makedirs("outputs", exist_ok=True)

            # Parse HTML for tables
            dfs = pd.read_html(io.StringIO(html_content))
            
            if not dfs:
                raise Exception("No tables found in HTML content")
            
            # Use the first table found
            df = dfs[0]
            
            # Save to CSV
            df.to_csv(output_path, index=False)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting HTML table to CSV: {str(e)}")
            raise Exception(f"Failed to convert HTML table to CSV: {str(e)}")

    @staticmethod
    def website_to_jpg(url: str, output_filename: str = None, width: int = 1920, height: int = 1080) -> str:
        """Convert website to JPG image."""
        # Determine filename
        if output_filename and output_filename.strip() and output_filename.lower() != "string":
            if not output_filename.lower().endswith('.jpg') and not output_filename.lower().endswith('.jpeg'):
                output_filename += '.jpg'
            filename = output_filename
        else:
            unique_id = str(uuid.uuid4())
            filename = f"website_to_jpg_{unique_id}.jpg"
        
        output_path = os.path.join("outputs", filename)
        os.makedirs("outputs", exist_ok=True)

        try:
            # Try using Selenium for real website rendering
            logger.info(f"Attempting to capture screenshot of {url} using Selenium...")
            
            chrome_options = Options()
            chrome_options.add_argument("--headless")
            chrome_options.add_argument("--no-sandbox")
            chrome_options.add_argument("--disable-dev-shm-usage")
            # Initial size, will be updated
            chrome_options.add_argument(f"--window-size={width},{height}")
            chrome_options.add_argument("--hide-scrollbars")
            
            # Setup driver
            service = ChromeService(ChromeDriverManager().install())
            driver = webdriver.Chrome(service=service, options=chrome_options)
            
            try:
                driver.get(url)
                # Optional: Wait for body to be present or just sleep briefly
                import time
                time.sleep(2) 
                
                # Get full page height
                total_height = driver.execute_script("return document.body.parentNode.scrollHeight")
                
                # Resize window to full height
                driver.set_window_size(width, total_height)
                
                # Wait a bit for resize to take effect
                time.sleep(1)
                
                # Save screenshot
                driver.save_screenshot(output_path)
                logger.info(f"Screenshot saved to {output_path}")
                return output_path
            finally:
                driver.quit()

        except Exception as e:
            logger.warning(f"Selenium conversion failed: {str(e)}. Falling back to placeholder.")
            
            try:
                # Create a placeholder image since Selenium failed
                img = Image.new('RGB', (width, height), color=(255, 255, 255))
                d = ImageDraw.Draw(img)
                
                # Try to load a font, fallback to default
                try:
                    font = ImageFont.truetype("arial.ttf", 40)
                    small_font = ImageFont.truetype("arial.ttf", 20)
                except IOError:
                    font = ImageFont.load_default()
                    small_font = ImageFont.load_default()

                # Draw text
                text = f"Website Preview: {url}"
                d.text((50, 50), text, fill=(0, 0, 0), font=font)
                
                error_msg = f"Error: Could not render website. {str(e)}"
                d.text((50, 100), error_msg, fill=(255, 0, 0), font=small_font)
                
                note = "Note: Please ensure Google Chrome is installed on the server."
                d.text((50, 150), note, fill=(100, 100, 100), font=small_font)
                
                timestamp = f"Generated: {uuid.uuid4()}"
                d.text((50, 200), timestamp, fill=(100, 100, 100), font=small_font)

                # Save image
                img.save(output_path, "JPEG")
                return output_path
                
            except Exception as fallback_error:
                logger.error(f"Fallback image generation failed: {str(fallback_error)}")
                raise Exception(f"Failed to convert website to JPG: {str(e)}")
    
    @staticmethod
    def html_to_jpg(html_content: str, original_filename: str = None, output_filename: str = None, width: int = 1920, height: int = 1080) -> str:
        """Convert HTML content to JPG image."""
        # Determine filename
        if output_filename and output_filename.strip() and output_filename.lower() != "string":
            if not output_filename.lower().endswith('.jpg') and not output_filename.lower().endswith('.jpeg'):
                output_filename += '.jpg'
            filename = output_filename
        else:
            if original_filename:
                base_name = os.path.splitext(original_filename)[0]
                filename = f"{base_name}.jpg"
            else:
                unique_id = str(uuid.uuid4())
                filename = f"html_to_jpg_{unique_id}.jpg"
        
        output_path = os.path.join("outputs", filename)
        os.makedirs("outputs", exist_ok=True)

        # Create a temporary HTML file to render
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as temp_html:
            temp_html.write(html_content)
            temp_html_path = temp_html.name

        try:
            # Try using Selenium for real website rendering
            file_url = f"file:///{temp_html_path.replace(os.sep, '/')}"
            logger.info(f"Attempting to capture screenshot of HTML content using Selenium...")
            
            chrome_options = Options()
            chrome_options.add_argument("--headless")
            chrome_options.add_argument("--no-sandbox")
            chrome_options.add_argument("--disable-dev-shm-usage")
            # Initial size, will be updated
            chrome_options.add_argument(f"--window-size={width},{height}")
            chrome_options.add_argument("--hide-scrollbars")
            
            # Setup driver
            service = ChromeService(ChromeDriverManager().install())
            driver = webdriver.Chrome(service=service, options=chrome_options)
            
            try:
                driver.get(file_url)
                # Optional: Wait for body to be present or just sleep briefly
                import time
                time.sleep(2) 
                
                # Get full page height using a more robust method
                total_height = driver.execute_script("""
                    return Math.max(
                        document.body.scrollHeight, 
                        document.body.offsetHeight, 
                        document.documentElement.clientHeight, 
                        document.documentElement.scrollHeight, 
                        document.documentElement.offsetHeight
                    );
                """)
                
                # Add a buffer to ensure bottom content is not cut off
                total_height += 100
                
                # Resize window to full height
                driver.set_window_size(width, total_height)
                
                # Wait a bit for resize to take effect
                time.sleep(1)
                
                # Save screenshot
                driver.save_screenshot(output_path)
                logger.info(f"Screenshot saved to {output_path}")
                return output_path
            finally:
                driver.quit()

        except Exception as e:
            logger.warning(f"Selenium conversion failed: {str(e)}. Falling back to placeholder.")
            
            try:
                # Create a placeholder image since Selenium failed
                img = Image.new('RGB', (width, height), color=(255, 255, 255))
                d = ImageDraw.Draw(img)
                
                # Try to load a font, fallback to default
                try:
                    font = ImageFont.truetype("arial.ttf", 40)
                    small_font = ImageFont.truetype("arial.ttf", 20)
                except IOError:
                    font = ImageFont.load_default()
                    small_font = ImageFont.load_default()

                # Draw text
                text = "HTML Preview"
                d.text((50, 50), text, fill=(0, 0, 0), font=font)
                
                error_msg = f"Error: Could not render HTML. {str(e)}"
                d.text((50, 100), error_msg, fill=(255, 0, 0), font=small_font)
                
                note = "Note: Please ensure Google Chrome is installed on the server."
                d.text((50, 150), note, fill=(100, 100, 100), font=small_font)
                
                timestamp = f"Generated: {uuid.uuid4()}"
                d.text((50, 200), timestamp, fill=(100, 100, 100), font=small_font)

                # Save image
                img.save(output_path, "JPEG")
                return output_path
                
            except Exception as fallback_error:
                logger.error(f"Fallback image generation failed: {str(fallback_error)}")
                raise Exception(f"Failed to convert HTML to JPG: {str(e)}")
        finally:
            # Clean up temp file
            if os.path.exists(temp_html_path):
                os.remove(temp_html_path)
    
    @staticmethod
    def website_to_png(url: str, output_filename: str = None, width: int = 1920, height: int = 1080) -> str:
        """Convert website to PNG image."""
        # Determine filename
        if output_filename and output_filename.strip() and output_filename.lower() != "string":
            if not output_filename.lower().endswith('.png'):
                output_filename += '.png'
            filename = output_filename
        else:
            unique_id = str(uuid.uuid4())
            filename = f"website_to_png_{unique_id}.png"
        
        output_path = os.path.join("outputs", filename)
        os.makedirs("outputs", exist_ok=True)

        try:
            # Try using Selenium for real website rendering
            logger.info(f"Attempting to capture screenshot of {url} using Selenium...")
            
            chrome_options = Options()
            chrome_options.add_argument("--headless")
            chrome_options.add_argument("--no-sandbox")
            chrome_options.add_argument("--disable-dev-shm-usage")
            # Initial size, will be updated
            chrome_options.add_argument(f"--window-size={width},{height}")
            chrome_options.add_argument("--hide-scrollbars")
            
            # Setup driver
            service = ChromeService(ChromeDriverManager().install())
            driver = webdriver.Chrome(service=service, options=chrome_options)
            
            try:
                driver.get(url)
                # Optional: Wait for body to be present or just sleep briefly
                import time
                time.sleep(2) 
                
                # Get full page height
                total_height = driver.execute_script("return document.body.parentNode.scrollHeight")
                
                # Resize window to full height
                driver.set_window_size(width, total_height)
                
                # Wait a bit for resize to take effect
                time.sleep(1)
                
                # Save screenshot (default is PNG)
                driver.save_screenshot(output_path)
                logger.info(f"Screenshot saved to {output_path}")
                return output_path
            finally:
                driver.quit()

        except Exception as e:
            logger.warning(f"Selenium conversion failed: {str(e)}. Falling back to placeholder.")
            
            try:
                # Create a placeholder image since Selenium failed
                img = Image.new('RGB', (width, height), color=(255, 255, 255))
                d = ImageDraw.Draw(img)
                
                # Try to load a font, fallback to default
                try:
                    font = ImageFont.truetype("arial.ttf", 40)
                    small_font = ImageFont.truetype("arial.ttf", 20)
                except IOError:
                    font = ImageFont.load_default()
                    small_font = ImageFont.load_default()

                # Draw text
                text = f"Website Preview: {url}"
                d.text((50, 50), text, fill=(0, 0, 0), font=font)
                
                error_msg = f"Error: Could not render website. {str(e)}"
                d.text((50, 100), error_msg, fill=(255, 0, 0), font=small_font)
                
                note = "Note: Please ensure Google Chrome is installed on the server."
                d.text((50, 150), note, fill=(100, 100, 100), font=small_font)
                
                timestamp = f"Generated: {uuid.uuid4()}"
                d.text((50, 200), timestamp, fill=(100, 100, 100), font=small_font)

                # Save image
                img.save(output_path, "PNG")
                return output_path
                
            except Exception as fallback_error:
                logger.error(f"Fallback image generation failed: {str(fallback_error)}")
                raise Exception(f"Failed to convert website to PNG: {str(e)}")
    
    @staticmethod
    def html_to_png(html_content: str, original_filename: str = None, output_filename: str = None, width: int = 1920, height: int = 1080) -> str:
        """Convert HTML content to PNG image."""
        # Determine filename
        if output_filename and output_filename.strip() and output_filename.lower() != "string":
            if not output_filename.lower().endswith('.png'):
                output_filename += '.png'
            filename = output_filename
        else:
            if original_filename:
                base_name = os.path.splitext(original_filename)[0]
                filename = f"{base_name}.png"
            else:
                unique_id = str(uuid.uuid4())
                filename = f"html_to_png_{unique_id}.png"
        
        output_path = os.path.join("outputs", filename)
        os.makedirs("outputs", exist_ok=True)

        # Create a temporary HTML file to render
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as temp_html:
            temp_html.write(html_content)
            temp_html_path = temp_html.name

        try:
            # Try using Selenium for real website rendering
            file_url = f"file:///{temp_html_path.replace(os.sep, '/')}"
            logger.info(f"Attempting to capture screenshot of HTML content using Selenium...")
            
            chrome_options = Options()
            chrome_options.add_argument("--headless")
            chrome_options.add_argument("--no-sandbox")
            chrome_options.add_argument("--disable-dev-shm-usage")
            # Initial size, will be updated
            chrome_options.add_argument(f"--window-size={width},{height}")
            chrome_options.add_argument("--hide-scrollbars")
            
            # Setup driver
            service = ChromeService(ChromeDriverManager().install())
            driver = webdriver.Chrome(service=service, options=chrome_options)
            
            try:
                driver.get(file_url)
                # Optional: Wait for body to be present or just sleep briefly
                import time
                time.sleep(2) 
                
                # Get full page height using a more robust method
                total_height = driver.execute_script("""
                    return Math.max(
                        document.body.scrollHeight, 
                        document.body.offsetHeight, 
                        document.documentElement.clientHeight, 
                        document.documentElement.scrollHeight, 
                        document.documentElement.offsetHeight
                    );
                """)
                
                # Add a buffer to ensure bottom content is not cut off
                total_height += 100
                
                # Resize window to full height
                driver.set_window_size(width, total_height)
                
                # Wait a bit for resize to take effect
                time.sleep(1)
                
                # Save screenshot (default is PNG)
                driver.save_screenshot(output_path)
                logger.info(f"Screenshot saved to {output_path}")
                return output_path
            finally:
                driver.quit()

        except Exception as e:
            logger.warning(f"Selenium conversion failed: {str(e)}. Falling back to placeholder.")
            
            try:
                # Create a placeholder image since Selenium failed
                img = Image.new('RGB', (width, height), color=(255, 255, 255))
                d = ImageDraw.Draw(img)
                
                # Try to load a font, fallback to default
                try:
                    font = ImageFont.truetype("arial.ttf", 40)
                    small_font = ImageFont.truetype("arial.ttf", 20)
                except IOError:
                    font = ImageFont.load_default()
                    small_font = ImageFont.load_default()

                # Draw text
                text = "HTML Preview"
                d.text((50, 50), text, fill=(0, 0, 0), font=font)
                
                error_msg = f"Error: Could not render HTML. {str(e)}"
                d.text((50, 100), error_msg, fill=(255, 0, 0), font=small_font)
                
                note = "Note: Please ensure Google Chrome is installed on the server."
                d.text((50, 150), note, fill=(100, 100, 100), font=small_font)
                
                timestamp = f"Generated: {uuid.uuid4()}"
                d.text((50, 200), timestamp, fill=(100, 100, 100), font=small_font)

                # Save image
                img.save(output_path, "PNG")
                return output_path
                
            except Exception as fallback_error:
                logger.error(f"Fallback image generation failed: {str(fallback_error)}")
                raise Exception(f"Failed to convert HTML to PNG: {str(e)}")
        finally:
            # Clean up temp file
            if os.path.exists(temp_html_path):
                os.remove(temp_html_path)
    

    
    @staticmethod
    def excel_to_html(file_content: bytes) -> str:
        """Convert Excel file to HTML."""
        try:
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted Excel</title></head><body>"
            
            for sheet_name, data in df.items():
                html_content += f"<h2>Sheet: {sheet_name}</h2>"
                html_content += data.to_html(index=False, classes='table table-striped')
                html_content += "<br>"
            
            html_content += "</body></html>"
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([excel_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting Excel to HTML: {str(e)}")
            raise Exception(f"Failed to convert Excel to HTML: {str(e)}")
    
    @staticmethod
    def pdf_to_html(file_content: bytes) -> str:
        """Convert PDF to HTML."""
        try:
            # Create temporary file for PDF
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_file:
                pdf_file.write(file_content)
                pdf_file_path = pdf_file.name
            
            # Read PDF
            doc = fitz.open(pdf_file_path)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted PDF</title></head><body>"
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    html_content += f"<div class='page'><h3>Page {page_num + 1}</h3>"
                    html_content += f"<p>{text}</p></div>"
            
            html_content += "</body></html>"
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([pdf_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting PDF to HTML: {str(e)}")
            raise Exception(f"Failed to convert PDF to HTML: {str(e)}")
    
    @staticmethod
    def log_conversion(conversion_type: str, input_data: str, output_data: str, success: bool, error_message: str = None, user_id: int = None):
        """Log conversion operation to database."""
        try:
            RequestLoggingService.log_conversion_request(
                conversion_type=conversion_type,
                input_data=input_data,
                output_data=output_data,
                success=success,
                error_message=error_message,
                user_id=user_id
            )
        except Exception as e:
            logger.error(f"Error logging conversion: {str(e)}")
    
    @staticmethod
    def cleanup_temp_files(file_paths: list):
        """Clean up temporary files."""
        for file_path in file_paths:
            if file_path and os.path.exists(file_path):
                try:
                    os.unlink(file_path)
                except Exception as e:
                    logger.warning(f"Could not delete temporary file {file_path}: {str(e)}")
