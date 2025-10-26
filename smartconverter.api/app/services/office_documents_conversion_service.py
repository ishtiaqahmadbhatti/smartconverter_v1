"""
Office Documents Conversion Service

This service handles various office document conversion operations including Excel, Word, PowerPoint,
OpenOffice, PDF, CSV, JSON, XML, BSON, and SRT files.
"""

import os
import json
import tempfile
import logging
import csv
import xml.etree.ElementTree as ET
from typing import Dict, Any, Optional, List
from pathlib import Path
import base64
import io
import re

# Document processing libraries
import pandas as pd
import docx
from pptx import Presentation
import fitz  # PyMuPDF
from bs4 import BeautifulSoup
import openpyxl
from openpyxl import Workbook
import xlsxwriter

# Database logging
from app.services.request_logging_service import RequestLoggingService

logger = logging.getLogger(__name__)


class OfficeDocumentsConversionService:
    """Service for office documents conversion operations."""
    
    # PDF Conversions
    @staticmethod
    def pdf_to_csv(file_content: bytes) -> str:
        """Convert PDF to CSV."""
        try:
            # Create temporary file for PDF
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_file:
                pdf_file.write(file_content)
                pdf_file_path = pdf_file.name
            
            # Read PDF
            doc = fitz.open(pdf_file_path)
            
            # Extract text and convert to CSV
            csv_content = "Page,Content\n"
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    # Clean text for CSV
                    clean_text = text.replace('"', '""').replace('\n', ' ')
                    csv_content += f"{page_num + 1},\"{clean_text}\"\n"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([pdf_file_path])
            
            return csv_content.strip()
            
        except Exception as e:
            logger.error(f"Error converting PDF to CSV: {str(e)}")
            raise Exception(f"Failed to convert PDF to CSV: {str(e)}")
    
    @staticmethod
    def pdf_to_excel(file_content: bytes) -> str:
        """Convert PDF to Excel file."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"pdf_to_excel_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for PDF
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_file:
                pdf_file.write(file_content)
                pdf_file_path = pdf_file.name
            
            # Read PDF
            doc = fitz.open(pdf_file_path)
            
            # Extract text and create Excel
            wb = Workbook()
            ws = wb.active
            ws.title = "PDF Content"
            
            # Add headers
            ws['A1'] = "Page"
            ws['B1'] = "Content"
            
            row = 2
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    ws.cell(row=row, column=1, value=page_num + 1)
                    ws.cell(row=row, column=2, value=text)
                    row += 1
            
            # Save Excel file
            wb.save(output_path)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([pdf_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting PDF to Excel: {str(e)}")
            raise Exception(f"Failed to convert PDF to Excel: {str(e)}")
    
    @staticmethod
    def pdf_to_word(file_content: bytes) -> str:
        """Convert PDF to Word document."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"pdf_to_word_{unique_id}.docx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for PDF
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_file:
                pdf_file.write(file_content)
                pdf_file_path = pdf_file.name
            
            # Read PDF
            doc = fitz.open(pdf_file_path)
            
            # Create Word document
            word_doc = docx.Document()
            
            # Extract text and add to Word
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    # Add page header
                    word_doc.add_heading(f"Page {page_num + 1}", level=1)
                    # Add content
                    word_doc.add_paragraph(text)
            
            # Save Word document
            word_doc.save(output_path)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([pdf_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting PDF to Word: {str(e)}")
            raise Exception(f"Failed to convert PDF to Word: {str(e)}")
    
    # Word Conversions
    @staticmethod
    def word_to_pdf(file_content: bytes) -> str:
        """Convert Word document to PDF."""
        try:
            import uuid
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"word_to_pdf_{unique_id}.pdf"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for Word
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as word_file:
                word_file.write(file_content)
                word_file_path = word_file.name
            
            # Read Word document
            doc = docx.Document(word_file_path)
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Process paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    para = Paragraph(paragraph.text, styles['Normal'])
                    story.append(para)
                    story.append(Spacer(1, 12))
            
            # Process tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        para = Paragraph(row_text, styles['Normal'])
                        story.append(para)
                        story.append(Spacer(1, 6))
            
            pdf_doc.build(story)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([word_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Word to PDF: {str(e)}")
            raise Exception(f"Failed to convert Word to PDF: {str(e)}")
    
    @staticmethod
    def word_to_html(file_content: bytes) -> str:
        """Convert Word document to HTML."""
        try:
            # Create temporary file for Word
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as word_file:
                word_file.write(file_content)
                word_file_path = word_file.name
            
            # Read Word document
            doc = docx.Document(word_file_path)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted Document</title></head><body>"
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Basic formatting
                    if paragraph.style.name.startswith('Heading'):
                        level = paragraph.style.name.split()[-1] if paragraph.style.name.split()[-1].isdigit() else '1'
                        html_content += f"<h{level}>{paragraph.text}</h{level}>"
                    else:
                        html_content += f"<p>{paragraph.text}</p>"
            
            # Add tables
            for table in doc.tables:
                html_content += "<table border='1'>"
                for row in table.rows:
                    html_content += "<tr>"
                    for cell in row.cells:
                        html_content += f"<td>{cell.text}</td>"
                    html_content += "</tr>"
                html_content += "</table>"
            
            html_content += "</body></html>"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([word_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting Word to HTML: {str(e)}")
            raise Exception(f"Failed to convert Word to HTML: {str(e)}")
    
    @staticmethod
    def word_to_text(file_content: bytes) -> str:
        """Convert Word document to text."""
        try:
            # Create temporary file for Word
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as word_file:
                word_file.write(file_content)
                word_file_path = word_file.name
            
            # Read Word document
            doc = docx.Document(word_file_path)
            
            # Extract text
            text_content = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"
            
            # Add tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_content += row_text + "\n"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([word_file_path])
            
            return text_content.strip()
            
        except Exception as e:
            logger.error(f"Error converting Word to text: {str(e)}")
            raise Exception(f"Failed to convert Word to text: {str(e)}")
    
    # PowerPoint Conversions
    @staticmethod
    def powerpoint_to_pdf(file_content: bytes) -> str:
        """Convert PowerPoint presentation to PDF."""
        try:
            import uuid
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"powerpoint_to_pdf_{unique_id}.pdf"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for PowerPoint
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as ppt_file:
                ppt_file.write(file_content)
                ppt_file_path = ppt_file.name
            
            # Read PowerPoint presentation
            prs = Presentation(ppt_file_path)
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for i, slide in enumerate(prs.slides, 1):
                story.append(Paragraph(f"Slide {i}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        para = Paragraph(shape.text, styles['Normal'])
                        story.append(para)
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 20))
            
            pdf_doc.build(story)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([ppt_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting PowerPoint to PDF: {str(e)}")
            raise Exception(f"Failed to convert PowerPoint to PDF: {str(e)}")
    
    @staticmethod
    def powerpoint_to_html(file_content: bytes) -> str:
        """Convert PowerPoint presentation to HTML."""
        try:
            # Create temporary file for PowerPoint
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as ppt_file:
                ppt_file.write(file_content)
                ppt_file_path = ppt_file.name
            
            # Read PowerPoint presentation
            prs = Presentation(ppt_file_path)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted Presentation</title></head><body>"
            
            for i, slide in enumerate(prs.slides, 1):
                html_content += f"<div class='slide'><h2>Slide {i}</h2>"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        html_content += f"<p>{shape.text}</p>"
                
                html_content += "</div><hr>"
            
            html_content += "</body></html>"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([ppt_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting PowerPoint to HTML: {str(e)}")
            raise Exception(f"Failed to convert PowerPoint to HTML: {str(e)}")
    
    @staticmethod
    def powerpoint_to_text(file_content: bytes) -> str:
        """Convert PowerPoint presentation to text."""
        try:
            # Create temporary file for PowerPoint
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as ppt_file:
                ppt_file.write(file_content)
                ppt_file_path = ppt_file.name
            
            # Read PowerPoint presentation
            prs = Presentation(ppt_file_path)
            
            # Extract text
            text_content = ""
            
            for i, slide in enumerate(prs.slides, 1):
                text_content += f"Slide {i}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                
                text_content += "\n"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([ppt_file_path])
            
            return text_content.strip()
            
        except Exception as e:
            logger.error(f"Error converting PowerPoint to text: {str(e)}")
            raise Exception(f"Failed to convert PowerPoint to text: {str(e)}")
    
    # Excel Conversions
    @staticmethod
    def excel_to_pdf(file_content: bytes) -> str:
        """Convert Excel file to PDF."""
        try:
            import uuid
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"excel_to_pdf_{unique_id}.pdf"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_name, data in df.items():
                story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                # Convert DataFrame to table
                table_data = [data.columns.tolist()] + data.values.tolist()
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 14),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(table)
                story.append(Spacer(1, 20))
            
            pdf_doc.build(story)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([excel_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Excel to PDF: {str(e)}")
            raise Exception(f"Failed to convert Excel to PDF: {str(e)}")
    
    @staticmethod
    def excel_to_xps(file_content: bytes) -> str:
        """Convert Excel file to XPS."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"excel_to_xps_{unique_id}.xps"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file and convert to XPS (simplified - save as text for now)
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            xps_content = "<?xml version='1.0' encoding='UTF-8'?>\n<XPSDocument>\n"
            
            for sheet_name, data in df.items():
                xps_content += f"  <Sheet name='{sheet_name}'>\n"
                for index, row in data.iterrows():
                    xps_content += f"    <Row index='{index}'>\n"
                    for column, value in row.items():
                        xps_content += f"      <Cell column='{column}'>{value}</Cell>\n"
                    xps_content += f"    </Row>\n"
                xps_content += f"  </Sheet>\n"
            
            xps_content += "</XPSDocument>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(xps_content)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([excel_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Excel to XPS: {str(e)}")
            raise Exception(f"Failed to convert Excel to XPS: {str(e)}")
    
    @staticmethod
    def excel_to_html(file_content: bytes) -> str:
        """Convert Excel file to HTML."""
        try:
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted Excel</title></head><body>"
            
            for sheet_name, data in df.items():
                html_content += f"<h2>Sheet: {sheet_name}</h2>"
                html_content += data.to_html(index=False, classes='table table-striped')
                html_content += "<br>"
            
            html_content += "</body></html>"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([excel_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting Excel to HTML: {str(e)}")
            raise Exception(f"Failed to convert Excel to HTML: {str(e)}")
    
    @staticmethod
    def excel_to_csv(file_content: bytes) -> str:
        """Convert Excel file to CSV."""
        try:
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            csv_content = ""
            
            for sheet_name, data in df.items():
                csv_content += f"# Sheet: {sheet_name}\n"
                csv_content += data.to_csv(index=False)
                csv_content += "\n\n"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([excel_file_path])
            
            return csv_content.strip()
            
        except Exception as e:
            logger.error(f"Error converting Excel to CSV: {str(e)}")
            raise Exception(f"Failed to convert Excel to CSV: {str(e)}")
    
    @staticmethod
    def excel_to_ods(file_content: bytes) -> str:
        """Convert Excel file to OpenOffice Calc ODS."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"excel_to_ods_{unique_id}.ods"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            # Save as ODS
            with pd.ExcelWriter(output_path, engine='odf') as writer:
                for sheet_name, data in df.items():
                    data.to_excel(writer, sheet_name=sheet_name, index=False)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([excel_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Excel to ODS: {str(e)}")
            raise Exception(f"Failed to convert Excel to ODS: {str(e)}")
    
    # OpenOffice Conversions
    @staticmethod
    def ods_to_csv(file_content: bytes) -> str:
        """Convert OpenOffice Calc ODS file to CSV."""
        try:
            # Create temporary file for ODS
            with tempfile.NamedTemporaryFile(suffix='.ods', delete=False) as ods_file:
                ods_file.write(file_content)
                ods_file_path = ods_file.name
            
            # Read ODS file
            df = pd.read_excel(ods_file_path, sheet_name=None, engine='odf')
            
            csv_content = ""
            
            for sheet_name, data in df.items():
                csv_content += f"# Sheet: {sheet_name}\n"
                csv_content += data.to_csv(index=False)
                csv_content += "\n\n"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([ods_file_path])
            
            return csv_content.strip()
            
        except Exception as e:
            logger.error(f"Error converting ODS to CSV: {str(e)}")
            raise Exception(f"Failed to convert ODS to CSV: {str(e)}")
    
    @staticmethod
    def ods_to_pdf(file_content: bytes) -> str:
        """Convert OpenOffice Calc ODS file to PDF."""
        try:
            import uuid
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"ods_to_pdf_{unique_id}.pdf"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for ODS
            with tempfile.NamedTemporaryFile(suffix='.ods', delete=False) as ods_file:
                ods_file.write(file_content)
                ods_file_path = ods_file.name
            
            # Read ODS file
            df = pd.read_excel(ods_file_path, sheet_name=None, engine='odf')
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_name, data in df.items():
                story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                # Convert DataFrame to table
                table_data = [data.columns.tolist()] + data.values.tolist()
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 14),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(table)
                story.append(Spacer(1, 20))
            
            pdf_doc.build(story)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([ods_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting ODS to PDF: {str(e)}")
            raise Exception(f"Failed to convert ODS to PDF: {str(e)}")
    
    @staticmethod
    def ods_to_excel(file_content: bytes) -> str:
        """Convert OpenOffice Calc ODS file to Excel."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"ods_to_excel_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Create temporary file for ODS
            with tempfile.NamedTemporaryFile(suffix='.ods', delete=False) as ods_file:
                ods_file.write(file_content)
                ods_file_path = ods_file.name
            
            # Read ODS file
            df = pd.read_excel(ods_file_path, sheet_name=None, engine='odf')
            
            # Save as Excel
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                for sheet_name, data in df.items():
                    data.to_excel(writer, sheet_name=sheet_name, index=False)
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([ods_file_path])
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting ODS to Excel: {str(e)}")
            raise Exception(f"Failed to convert ODS to Excel: {str(e)}")
    
    # CSV Conversions
    @staticmethod
    def csv_to_excel(csv_content: str) -> str:
        """Convert CSV to Excel file."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"csv_to_excel_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Read CSV content
            from io import StringIO
            df = pd.read_csv(StringIO(csv_content))
            
            # Save as Excel
            df.to_excel(output_path, index=False)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting CSV to Excel: {str(e)}")
            raise Exception(f"Failed to convert CSV to Excel: {str(e)}")
    
    # XML Conversions
    @staticmethod
    def excel_to_xml(file_content: bytes, root_name: str = "data", record_name: str = "record") -> str:
        """Convert Excel file to XML."""
        try:
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            xml_content = f'<?xml version="1.0" encoding="UTF-8"?>\n<{root_name}>\n'
            
            for sheet_name, data in df.items():
                xml_content += f'  <sheet name="{sheet_name}">\n'
                for index, row in data.iterrows():
                    xml_content += f'    <{record_name} id="{index}">\n'
                    for column, value in row.items():
                        clean_column = column.replace(' ', '_').replace('-', '_').replace('(', '').replace(')', '')
                        escaped_value = OfficeDocumentsConversionService._escape_xml(str(value))
                        xml_content += f'      <{clean_column}>{escaped_value}</{clean_column}>\n'
                    xml_content += f'    </{record_name}>\n'
                xml_content += f'  </sheet>\n'
            
            xml_content += f'</{root_name}>'
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([excel_file_path])
            
            return xml_content
            
        except Exception as e:
            logger.error(f"Error converting Excel to XML: {str(e)}")
            raise Exception(f"Failed to convert Excel to XML: {str(e)}")
    
    @staticmethod
    def xml_to_csv(xml_content: str) -> str:
        """Convert XML to CSV."""
        try:
            # Parse XML
            root = ET.fromstring(xml_content)
            
            # Extract data
            records = []
            for record in root.findall('.//record'):
                record_data = {}
                for child in record:
                    record_data[child.tag] = child.text or ''
                records.append(record_data)
            
            if not records:
                raise Exception("No records found in XML")
            
            # Convert to CSV
            import io
            output = io.StringIO()
            
            # Get all unique keys
            all_keys = set()
            for record in records:
                all_keys.update(record.keys())
            
            # Write CSV
            writer = csv.DictWriter(output, fieldnames=sorted(all_keys))
            writer.writeheader()
            writer.writerows(records)
            
            return output.getvalue()
            
        except Exception as e:
            logger.error(f"Error converting XML to CSV: {str(e)}")
            raise Exception(f"Failed to convert XML to CSV: {str(e)}")
    
    @staticmethod
    def xml_to_excel(xml_content: str) -> str:
        """Convert XML to Excel file."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"xml_to_excel_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Parse XML
            root = ET.fromstring(xml_content)
            
            # Extract data
            records = []
            for record in root.findall('.//record'):
                record_data = {}
                for child in record:
                    record_data[child.tag] = child.text or ''
                records.append(record_data)
            
            if not records:
                raise Exception("No records found in XML")
            
            # Convert to DataFrame and save as Excel
            df = pd.DataFrame(records)
            df.to_excel(output_path, index=False)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting XML to Excel: {str(e)}")
            raise Exception(f"Failed to convert XML to Excel: {str(e)}")
    
    @staticmethod
    def excel_xml_to_xlsx(file_content: bytes) -> str:
        """Convert Excel XML to Excel XLSX file."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"excel_xml_to_xlsx_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Parse XML content
            xml_content = file_content.decode('utf-8')
            root = ET.fromstring(xml_content)
            
            # Create new Excel workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "Sheet1"
            
            # Extract data from XML
            records = []
            for record in root.findall('.//record'):
                record_data = {}
                for child in record:
                    record_data[child.tag] = child.text or ''
                records.append(record_data)
            
            if records:
                # Write headers
                headers = list(records[0].keys())
                for col, header in enumerate(headers, 1):
                    ws.cell(row=1, column=col, value=header)
                
                # Write data
                for row, record in enumerate(records, 2):
                    for col, header in enumerate(headers, 1):
                        ws.cell(row=row, column=col, value=record.get(header, ''))
            
            # Save workbook
            wb.save(output_path)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Excel XML to XLSX: {str(e)}")
            raise Exception(f"Failed to convert Excel XML to XLSX: {str(e)}")
    
    # JSON Conversions
    @staticmethod
    def json_to_excel(json_data: dict) -> str:
        """Convert JSON to Excel file."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"json_to_excel_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Convert JSON to DataFrame
            if isinstance(json_data, list):
                df = pd.DataFrame(json_data)
            elif isinstance(json_data, dict):
                # Flatten nested dictionaries
                flattened = OfficeDocumentsConversionService._flatten_dict(json_data)
                df = pd.DataFrame([flattened])
            else:
                raise Exception("Invalid JSON format")
            
            # Save as Excel
            df.to_excel(output_path, index=False)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting JSON to Excel: {str(e)}")
            raise Exception(f"Failed to convert JSON to Excel: {str(e)}")
    
    @staticmethod
    def excel_to_json(file_content: bytes) -> str:
        """Convert Excel file to JSON."""
        try:
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            # Convert to JSON
            json_data = {}
            for sheet_name, data in df.items():
                json_data[sheet_name] = data.to_dict('records')
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([excel_file_path])
            
            return json.dumps(json_data, indent=2)
            
        except Exception as e:
            logger.error(f"Error converting Excel to JSON: {str(e)}")
            raise Exception(f"Failed to convert Excel to JSON: {str(e)}")
    
    @staticmethod
    def json_objects_to_excel(json_objects: list) -> str:
        """Convert JSON objects to Excel file."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"json_objects_to_excel_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Convert JSON objects to DataFrame
            df = pd.DataFrame(json_objects)
            
            # Save as Excel
            df.to_excel(output_path, index=False)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting JSON objects to Excel: {str(e)}")
            raise Exception(f"Failed to convert JSON objects to Excel: {str(e)}")
    
    # BSON Conversions
    @staticmethod
    def bson_to_excel(bson_data: bytes) -> str:
        """Convert BSON to Excel file."""
        try:
            import uuid
            import bson
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"bson_to_excel_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Parse BSON
            bson_docs = bson.decode_all(bson_data)
            
            # Convert to DataFrame
            df = pd.DataFrame(bson_docs)
            
            # Save as Excel
            df.to_excel(output_path, index=False)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting BSON to Excel: {str(e)}")
            raise Exception(f"Failed to convert BSON to Excel: {str(e)}")
    
    # SRT Conversions
    @staticmethod
    def srt_to_excel(srt_content: str) -> str:
        """Convert SRT subtitle file to Excel."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"srt_to_excel_{unique_id}.xlsx"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Parse SRT content
            srt_entries = OfficeDocumentsConversionService._parse_srt(srt_content)
            
            # Convert to DataFrame
            df = pd.DataFrame(srt_entries)
            
            # Save as Excel
            df.to_excel(output_path, index=False)
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting SRT to Excel: {str(e)}")
            raise Exception(f"Failed to convert SRT to Excel: {str(e)}")
    
    @staticmethod
    def srt_to_xlsx(srt_content: str) -> str:
        """Convert SRT subtitle file to XLSX."""
        return OfficeDocumentsConversionService.srt_to_excel(srt_content)
    
    @staticmethod
    def srt_to_xls(srt_content: str) -> str:
        """Convert SRT subtitle file to XLS."""
        try:
            import uuid
            
            # Create unique filename
            unique_id = str(uuid.uuid4())
            filename = f"srt_to_xls_{unique_id}.xls"
            output_path = os.path.join("outputs", filename)
            
            # Ensure outputs directory exists
            os.makedirs("outputs", exist_ok=True)
            
            # Parse SRT content
            srt_entries = OfficeDocumentsConversionService._parse_srt(srt_content)
            
            # Convert to DataFrame
            df = pd.DataFrame(srt_entries)
            
            # Save as XLS (using xlsxwriter for compatibility)
            df.to_excel(output_path, index=False, engine='xlsxwriter')
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting SRT to XLS: {str(e)}")
            raise Exception(f"Failed to convert SRT to XLS: {str(e)}")
    
    @staticmethod
    def excel_to_srt(file_content: bytes) -> str:
        """Convert Excel file to SRT subtitle file."""
        try:
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            # Convert to SRT
            srt_content = ""
            
            for sheet_name, data in df.items():
                for index, row in data.iterrows():
                    srt_content += f"{index + 1}\n"
                    srt_content += f"{row.get('Start_Time', '00:00:00,000')} --> {row.get('End_Time', '00:00:00,000')}\n"
                    srt_content += f"{row.get('Text', '')}\n\n"
            
            # Cleanup
            OfficeDocumentsConversionService.cleanup_temp_files([excel_file_path])
            
            return srt_content.strip()
            
        except Exception as e:
            logger.error(f"Error converting Excel to SRT: {str(e)}")
            raise Exception(f"Failed to convert Excel to SRT: {str(e)}")
    
    @staticmethod
    def xlsx_to_srt(file_content: bytes) -> str:
        """Convert XLSX file to SRT subtitle file."""
        return OfficeDocumentsConversionService.excel_to_srt(file_content)
    
    @staticmethod
    def xls_to_srt(file_content: bytes) -> str:
        """Convert XLS file to SRT subtitle file."""
        return OfficeDocumentsConversionService.excel_to_srt(file_content)
    
    # Helper Methods
    @staticmethod
    def _escape_xml(text: str) -> str:
        """Escape XML special characters."""
        if not text:
            return ""
        
        # Replace XML special characters
        text = text.replace('&', '&amp;')
        text = text.replace('<', '&lt;')
        text = text.replace('>', '&gt;')
        text = text.replace('"', '&quot;')
        text = text.replace("'", '&apos;')
        
        return text
    
    @staticmethod
    def _flatten_dict(d: dict, parent_key: str = '', sep: str = '_') -> dict:
        """Flatten nested dictionary."""
        items = []
        for k, v in d.items():
            new_key = f"{parent_key}{sep}{k}" if parent_key else k
            if isinstance(v, dict):
                items.extend(OfficeDocumentsConversionService._flatten_dict(v, new_key, sep=sep).items())
            else:
                items.append((new_key, v))
        return dict(items)
    
    @staticmethod
    def _parse_srt(srt_content: str) -> list:
        """Parse SRT content into list of dictionaries."""
        entries = []
        blocks = srt_content.strip().split('\n\n')
        
        for block in blocks:
            lines = block.strip().split('\n')
            if len(lines) >= 3:
                try:
                    index = int(lines[0])
                    time_line = lines[1]
                    text = '\n'.join(lines[2:])
                    
                    # Parse time
                    if ' --> ' in time_line:
                        start_time, end_time = time_line.split(' --> ')
                        entries.append({
                            'index': index,
                            'start_time': start_time.strip(),
                            'end_time': end_time.strip(),
                            'text': text
                        })
                except (ValueError, IndexError):
                    continue
        
        return entries
    
    @staticmethod
    def log_conversion(conversion_type: str, input_data: str, output_data: str, success: bool, error_message: str = None, user_id: int = None):
        """Log conversion operation to database."""
        try:
            RequestLoggingService.log_conversion_request(
                conversion_type=conversion_type,
                input_data=input_data,
                output_data=output_data,
                success=success,
                error_message=error_message,
                user_id=user_id
            )
        except Exception as e:
            logger.error(f"Error logging conversion: {str(e)}")
    
    @staticmethod
    def cleanup_temp_files(file_paths: list):
        """Clean up temporary files."""
        for file_path in file_paths:
            if file_path and os.path.exists(file_path):
                try:
                    os.unlink(file_path)
                except Exception as e:
                    logger.warning(f"Could not delete temporary file {file_path}: {str(e)}")
