import os
import io
from typing import Optional, Dict, Any, List, Tuple
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import pysrt
import webvtt
from app.core.exceptions import FileProcessingError
from app.services.file_service import FileService


class TextConversionService:
    """Service for handling text extraction from various document formats."""
    
    # Supported input formats
    SUPPORTED_INPUT_FORMATS = {
        'DOCX', 'DOC', 'PPTX', 'PPT', 'PDF', 'SRT', 'VTT', 'TXT'
    }
    
    @staticmethod
    def word_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from Word document."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input Word file not found: {input_path}")
            
            # Load Word document
            doc = Document(input_path)
            
            # Extract text from all paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Word to text conversion failed: {str(e)}")
    
    @staticmethod
    def powerpoint_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input PowerPoint file not found: {input_path}")
            
            # Load PowerPoint presentation
            prs = Presentation(input_path)
            
            # Extract text from all slides
            text_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Extract text from tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.extend(slide_text)
                    text_content.append("")  # Add empty line between slides
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"PowerPoint to text conversion failed: {str(e)}")
    
    @staticmethod
    def pdf_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from PDF document."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input PDF file not found: {input_path}")
            
            # Open PDF document
            doc = fitz.open(input_path)
            
            # Extract text from all pages
            text_content = []
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_text = page.get_text()
                
                if page_text.strip():
                    text_content.append(f"--- Page {page_num + 1} ---")
                    text_content.append(page_text.strip())
                    text_content.append("")  # Add empty line between pages
            
            # Close document
            doc.close()
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"PDF to text conversion failed: {str(e)}")
    
    @staticmethod
    def srt_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from SRT subtitle file."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input SRT file not found: {input_path}")
            
            # Load SRT file
            subs = pysrt.open(input_path)
            
            # Extract text from all subtitles
            text_content = []
            for sub in subs:
                if sub.text.strip():
                    text_content.append(sub.text.strip())
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"SRT to text conversion failed: {str(e)}")
    
    @staticmethod
    def vtt_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from VTT subtitle file."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input VTT file not found: {input_path}")
            
            # Load VTT file
            vtt = webvtt.read(input_path)
            
            # Extract text from all captions
            text_content = []
            for caption in vtt:
                if caption.text.strip():
                    text_content.append(caption.text.strip())
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"VTT to text conversion failed: {str(e)}")
    
    @staticmethod
    def get_supported_formats() -> List[str]:
        """Get list of supported input formats."""
        return list(TextConversionService.SUPPORTED_INPUT_FORMATS)
    
    @staticmethod
    def cleanup_temp_files(*file_paths: str) -> None:
        """Clean up temporary files."""
        for file_path in file_paths:
            FileService.cleanup_file(file_path)
