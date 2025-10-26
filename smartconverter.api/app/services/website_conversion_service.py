"""
Website Conversion Service

This service handles various website and HTML conversion operations including:
- HTML to PDF
- Word to HTML
- PowerPoint to HTML
- Markdown to HTML
- Website/HTML to JPG/PNG
- HTML Table to CSV
- Excel to HTML
- PDF to HTML
"""

import os
import json
import tempfile
import logging
from typing import Dict, Any, Optional
from pathlib import Path
import base64
import io

# HTML/Web conversion libraries
# from weasyprint import HTML, CSS  # Commented out due to Windows compatibility issues
from markdown import markdown
from bs4 import BeautifulSoup
import requests
from PIL import Image
import pdfkit
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

# Document processing libraries
import docx
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook
import fitz  # PyMuPDF

# Database logging
from app.services.request_logging_service import RequestLoggingService

logger = logging.getLogger(__name__)


class WebsiteConversionService:
    """Service for website and HTML conversion operations."""
    
    @staticmethod
    def html_to_pdf(html_content: str, css_content: str = None) -> str:
        """Convert HTML content to PDF using pdfkit (wkhtmltopdf)."""
        try:
            # Create temporary file for HTML
            with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as html_file:
                html_file.write(html_content)
                html_file_path = html_file.name
            
            # Convert to PDF using pdfkit
            output_path = tempfile.mktemp(suffix='.pdf')
            
            # Configure pdfkit options
            options = {
                'page-size': 'A4',
                'margin-top': '0.75in',
                'margin-right': '0.75in',
                'margin-bottom': '0.75in',
                'margin-left': '0.75in',
                'encoding': "UTF-8",
                'no-outline': None
            }
            
            # Convert HTML to PDF
            pdfkit.from_file(html_file_path, output_path, options=options)
            
            # Read PDF content
            with open(output_path, 'rb') as pdf_file:
                pdf_content = pdf_file.read()
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([html_file_path, output_path])
            
            return base64.b64encode(pdf_content).decode('utf-8')
            
        except Exception as e:
            logger.error(f"Error converting HTML to PDF: {str(e)}")
            raise Exception(f"Failed to convert HTML to PDF: {str(e)}")
    
    @staticmethod
    def word_to_html(file_content: bytes) -> str:
        """Convert Word document to HTML."""
        try:
            # Create temporary file for Word document
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as word_file:
                word_file.write(file_content)
                word_file_path = word_file.name
            
            # Read Word document
            doc = docx.Document(word_file_path)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted Document</title></head><body>"
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Basic formatting
                    if paragraph.style.name.startswith('Heading'):
                        level = paragraph.style.name.split()[-1] if paragraph.style.name.split()[-1].isdigit() else '1'
                        html_content += f"<h{level}>{paragraph.text}</h{level}>"
                    else:
                        html_content += f"<p>{paragraph.text}</p>"
            
            # Add tables
            for table in doc.tables:
                html_content += "<table border='1'>"
                for row in table.rows:
                    html_content += "<tr>"
                    for cell in row.cells:
                        html_content += f"<td>{cell.text}</td>"
                    html_content += "</tr>"
                html_content += "</table>"
            
            html_content += "</body></html>"
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([word_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting Word to HTML: {str(e)}")
            raise Exception(f"Failed to convert Word to HTML: {str(e)}")
    
    @staticmethod
    def powerpoint_to_html(file_content: bytes) -> str:
        """Convert PowerPoint presentation to HTML."""
        try:
            # Create temporary file for PowerPoint
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as ppt_file:
                ppt_file.write(file_content)
                ppt_file_path = ppt_file.name
            
            # Read PowerPoint presentation
            prs = Presentation(ppt_file_path)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted Presentation</title></head><body>"
            
            for i, slide in enumerate(prs.slides, 1):
                html_content += f"<div class='slide'><h2>Slide {i}</h2>"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        html_content += f"<p>{shape.text}</p>"
                
                html_content += "</div><hr>"
            
            html_content += "</body></html>"
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([ppt_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting PowerPoint to HTML: {str(e)}")
            raise Exception(f"Failed to convert PowerPoint to HTML: {str(e)}")
    
    @staticmethod
    def markdown_to_html(markdown_content: str) -> str:
        """Convert Markdown content to HTML."""
        try:
            html_content = markdown(markdown_content, extensions=['tables', 'fenced_code', 'codehilite'])
            
            # Wrap in proper HTML structure
            full_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <title>Converted Markdown</title>
                <style>
                    body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                    code {{ background-color: #f4f4f4; padding: 2px 4px; border-radius: 3px; }}
                    pre {{ background-color: #f4f4f4; padding: 10px; border-radius: 5px; overflow-x: auto; }}
                    table {{ border-collapse: collapse; width: 100%; }}
                    th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                    th {{ background-color: #f2f2f2; }}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
            
            return full_html
            
        except Exception as e:
            logger.error(f"Error converting Markdown to HTML: {str(e)}")
            raise Exception(f"Failed to convert Markdown to HTML: {str(e)}")
    
    @staticmethod
    def website_to_jpg(url: str, width: int = 1920, height: int = 1080) -> str:
        """Convert website to JPG image."""
        try:
            # Setup Chrome options for headless browsing
            chrome_options = Options()
            chrome_options.add_argument('--headless')
            chrome_options.add_argument('--no-sandbox')
            chrome_options.add_argument('--disable-dev-shm-usage')
            chrome_options.add_argument(f'--window-size={width},{height}')
            
            # Initialize WebDriver
            driver = webdriver.Chrome(options=chrome_options)
            
            try:
                # Navigate to website
                driver.get(url)
                
                # Wait for page to load
                WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
                
                # Take screenshot
                screenshot = driver.get_screenshot_as_png()
                
                # Convert to JPG
                image = Image.open(io.BytesIO(screenshot))
                if image.mode in ('RGBA', 'LA', 'P'):
                    image = image.convert('RGB')
                
                # Save to temporary file
                output_path = tempfile.mktemp(suffix='.jpg')
                image.save(output_path, 'JPEG', quality=95)
                
                # Read and encode
                with open(output_path, 'rb') as img_file:
                    img_content = img_file.read()
                
                # Cleanup
                WebsiteConversionService.cleanup_temp_files([output_path])
                
                return base64.b64encode(img_content).decode('utf-8')
                
            finally:
                driver.quit()
                
        except Exception as e:
            logger.error(f"Error converting website to JPG: {str(e)}")
            raise Exception(f"Failed to convert website to JPG: {str(e)}")
    
    @staticmethod
    def html_to_jpg(html_content: str, width: int = 1920, height: int = 1080) -> str:
        """Convert HTML content to JPG image."""
        try:
            # Create temporary HTML file
            with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as html_file:
                html_file.write(html_content)
                html_file_path = html_file.name
            
            # Setup Chrome options
            chrome_options = Options()
            chrome_options.add_argument('--headless')
            chrome_options.add_argument('--no-sandbox')
            chrome_options.add_argument('--disable-dev-shm-usage')
            chrome_options.add_argument(f'--window-size={width},{height}')
            
            # Initialize WebDriver
            driver = webdriver.Chrome(options=chrome_options)
            
            try:
                # Load HTML file
                file_url = f"file://{os.path.abspath(html_file_path)}"
                driver.get(file_url)
                
                # Wait for page to load
                WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
                
                # Take screenshot
                screenshot = driver.get_screenshot_as_png()
                
                # Convert to JPG
                image = Image.open(io.BytesIO(screenshot))
                if image.mode in ('RGBA', 'LA', 'P'):
                    image = image.convert('RGB')
                
                # Save to temporary file
                output_path = tempfile.mktemp(suffix='.jpg')
                image.save(output_path, 'JPEG', quality=95)
                
                # Read and encode
                with open(output_path, 'rb') as img_file:
                    img_content = img_file.read()
                
                # Cleanup
                WebsiteConversionService.cleanup_temp_files([html_file_path, output_path])
                
                return base64.b64encode(img_content).decode('utf-8')
                
            finally:
                driver.quit()
                
        except Exception as e:
            logger.error(f"Error converting HTML to JPG: {str(e)}")
            raise Exception(f"Failed to convert HTML to JPG: {str(e)}")
    
    @staticmethod
    def website_to_png(url: str, width: int = 1920, height: int = 1080) -> str:
        """Convert website to PNG image."""
        try:
            # Setup Chrome options for headless browsing
            chrome_options = Options()
            chrome_options.add_argument('--headless')
            chrome_options.add_argument('--no-sandbox')
            chrome_options.add_argument('--disable-dev-shm-usage')
            chrome_options.add_argument(f'--window-size={width},{height}')
            
            # Initialize WebDriver
            driver = webdriver.Chrome(options=chrome_options)
            
            try:
                # Navigate to website
                driver.get(url)
                
                # Wait for page to load
                WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
                
                # Take screenshot
                screenshot = driver.get_screenshot_as_png()
                
                # Save to temporary file
                output_path = tempfile.mktemp(suffix='.png')
                with open(output_path, 'wb') as img_file:
                    img_file.write(screenshot)
                
                # Read and encode
                with open(output_path, 'rb') as img_file:
                    img_content = img_file.read()
                
                # Cleanup
                WebsiteConversionService.cleanup_temp_files([output_path])
                
                return base64.b64encode(img_content).decode('utf-8')
                
            finally:
                driver.quit()
                
        except Exception as e:
            logger.error(f"Error converting website to PNG: {str(e)}")
            raise Exception(f"Failed to convert website to PNG: {str(e)}")
    
    @staticmethod
    def html_to_png(html_content: str, width: int = 1920, height: int = 1080) -> str:
        """Convert HTML content to PNG image."""
        try:
            # Create temporary HTML file
            with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as html_file:
                html_file.write(html_content)
                html_file_path = html_file.name
            
            # Setup Chrome options
            chrome_options = Options()
            chrome_options.add_argument('--headless')
            chrome_options.add_argument('--no-sandbox')
            chrome_options.add_argument('--disable-dev-shm-usage')
            chrome_options.add_argument(f'--window-size={width},{height}')
            
            # Initialize WebDriver
            driver = webdriver.Chrome(options=chrome_options)
            
            try:
                # Load HTML file
                file_url = f"file://{os.path.abspath(html_file_path)}"
                driver.get(file_url)
                
                # Wait for page to load
                WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
                
                # Take screenshot
                screenshot = driver.get_screenshot_as_png()
                
                # Save to temporary file
                output_path = tempfile.mktemp(suffix='.png')
                with open(output_path, 'wb') as img_file:
                    img_file.write(screenshot)
                
                # Read and encode
                with open(output_path, 'rb') as img_file:
                    img_content = img_file.read()
                
                # Cleanup
                WebsiteConversionService.cleanup_temp_files([html_file_path, output_path])
                
                return base64.b64encode(img_content).decode('utf-8')
                
            finally:
                driver.quit()
                
        except Exception as e:
            logger.error(f"Error converting HTML to PNG: {str(e)}")
            raise Exception(f"Failed to convert HTML to PNG: {str(e)}")
    
    @staticmethod
    def html_table_to_csv(html_content: str) -> str:
        """Convert HTML table to CSV."""
        try:
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Find tables
            tables = soup.find_all('table')
            if not tables:
                raise Exception("No tables found in HTML content")
            
            csv_content = ""
            
            for table in tables:
                rows = table.find_all('tr')
                for row in rows:
                    cells = row.find_all(['td', 'th'])
                    row_data = [cell.get_text(strip=True) for cell in cells]
                    csv_content += ','.join(f'"{cell}"' for cell in row_data) + '\n'
                csv_content += '\n'  # Separate tables
            
            return csv_content.strip()
            
        except Exception as e:
            logger.error(f"Error converting HTML table to CSV: {str(e)}")
            raise Exception(f"Failed to convert HTML table to CSV: {str(e)}")
    
    @staticmethod
    def excel_to_html(file_content: bytes) -> str:
        """Convert Excel file to HTML."""
        try:
            # Create temporary file for Excel
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as excel_file:
                excel_file.write(file_content)
                excel_file_path = excel_file.name
            
            # Read Excel file
            df = pd.read_excel(excel_file_path, sheet_name=None)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted Excel</title></head><body>"
            
            for sheet_name, data in df.items():
                html_content += f"<h2>Sheet: {sheet_name}</h2>"
                html_content += data.to_html(index=False, classes='table table-striped')
                html_content += "<br>"
            
            html_content += "</body></html>"
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([excel_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting Excel to HTML: {str(e)}")
            raise Exception(f"Failed to convert Excel to HTML: {str(e)}")
    
    @staticmethod
    def pdf_to_html(file_content: bytes) -> str:
        """Convert PDF to HTML."""
        try:
            # Create temporary file for PDF
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_file:
                pdf_file.write(file_content)
                pdf_file_path = pdf_file.name
            
            # Read PDF
            doc = fitz.open(pdf_file_path)
            
            # Convert to HTML
            html_content = "<html><head><title>Converted PDF</title></head><body>"
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    html_content += f"<div class='page'><h3>Page {page_num + 1}</h3>"
                    html_content += f"<p>{text}</p></div>"
            
            html_content += "</body></html>"
            
            # Cleanup
            WebsiteConversionService.cleanup_temp_files([pdf_file_path])
            
            return html_content
            
        except Exception as e:
            logger.error(f"Error converting PDF to HTML: {str(e)}")
            raise Exception(f"Failed to convert PDF to HTML: {str(e)}")
    
    @staticmethod
    def log_conversion(conversion_type: str, input_data: str, output_data: str, success: bool, error_message: str = None, user_id: int = None):
        """Log conversion operation to database."""
        try:
            RequestLoggingService.log_conversion_request(
                conversion_type=conversion_type,
                input_data=input_data,
                output_data=output_data,
                success=success,
                error_message=error_message,
                user_id=user_id
            )
        except Exception as e:
            logger.error(f"Error logging conversion: {str(e)}")
    
    @staticmethod
    def cleanup_temp_files(file_paths: list):
        """Clean up temporary files."""
        for file_path in file_paths:
            if file_path and os.path.exists(file_path):
                try:
                    os.unlink(file_path)
                except Exception as e:
                    logger.warning(f"Could not delete temporary file {file_path}: {str(e)}")
