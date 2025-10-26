#!/usr/bin/env python3
"""
Test script for text conversion functionality.
"""

import requests
import os
import json
import tempfile
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Test configuration
BASE_URL = "http://localhost:8000/api/v1/textconversiontools"
TEST_DOCX_PATH = "test_document.docx"
TEST_PPTX_PATH = "test_presentation.pptx"
TEST_PDF_PATH = "test_document.pdf"
TEST_SRT_PATH = "test_subtitles.srt"
TEST_VTT_PATH = "test_subtitles.vtt"

def create_test_docx():
    """Create a test Word document."""
    doc = Document()
    
    # Add title
    title = doc.add_heading('Test Document', 0)
    
    # Add paragraph
    doc.add_paragraph('This is a test paragraph in the Word document.')
    
    # Add another paragraph
    doc.add_paragraph('This is the second paragraph with some content.')
    
    # Add a table
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = 'Header 1'
    table.cell(0, 1).text = 'Header 2'
    table.cell(1, 0).text = 'Data 1'
    table.cell(1, 1).text = 'Data 2'
    
    # Add another paragraph
    doc.add_paragraph('This is the final paragraph of the test document.')
    
    doc.save(TEST_DOCX_PATH)
    print(f"Created test Word document: {TEST_DOCX_PATH}")

def create_test_pptx():
    """Create a test PowerPoint presentation."""
    prs = Presentation()
    
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "Test Presentation"
    subtitle1.text = "This is the first slide"
    
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    title2.text = "Slide 2"
    content2.text = "This is the second slide with some content."
    
    # Slide 3
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    title3.text = "Final Slide"
    content3.text = "This is the final slide of the presentation."
    
    prs.save(TEST_PPTX_PATH)
    print(f"Created test PowerPoint presentation: {TEST_PPTX_PATH}")

def create_test_pdf():
    """Create a test PDF document."""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        
        c = canvas.Canvas(TEST_PDF_PATH, pagesize=letter)
        
        # Page 1
        c.drawString(100, 750, "Test PDF Document")
        c.drawString(100, 700, "This is a test PDF document.")
        c.drawString(100, 650, "It contains multiple lines of text.")
        c.drawString(100, 600, "This is the first page.")
        c.showPage()
        
        # Page 2
        c.drawString(100, 750, "Page 2")
        c.drawString(100, 700, "This is the second page of the PDF.")
        c.drawString(100, 650, "It contains additional content.")
        c.drawString(100, 600, "This is the final page.")
        c.save()
        
        print(f"Created test PDF document: {TEST_PDF_PATH}")
    except ImportError:
        print("⚠️ ReportLab not available - skipping PDF creation")
        return False
    return True

def create_test_srt():
    """Create a test SRT file."""
    srt_content = """1
00:00:01,000 --> 00:00:04,000
Hello, this is a test subtitle.

2
00:00:05,000 --> 00:00:08,000
This is the second subtitle line.

3
00:00:10,000 --> 00:00:13,000
And this is the third subtitle.
"""
    
    with open(TEST_SRT_PATH, 'w', encoding='utf-8') as f:
        f.write(srt_content)
    
    print(f"Created test SRT file: {TEST_SRT_PATH}")

def create_test_vtt():
    """Create a test VTT file."""
    vtt_content = """WEBVTT

00:00:01.000 --> 00:00:04.000
Hello, this is a test subtitle.

00:00:05.000 --> 00:00:08.000
This is the second subtitle line.

00:00:10.000 --> 00:00:13.000
And this is the third subtitle.
"""
    
    with open(TEST_VTT_PATH, 'w', encoding='utf-8') as f:
        f.write(vtt_content)
    
    print(f"Created test VTT file: {TEST_VTT_PATH}")

def test_supported_formats():
    """Test getting supported formats."""
    try:
        response = requests.get(f"{BASE_URL}/supported-formats")
        if response.status_code == 200:
            data = response.json()
            print("✅ Supported formats test passed")
            print(f"Supported formats: {json.dumps(data, indent=2)}")
        else:
            print(f"❌ Supported formats test failed: {response.status_code}")
    except Exception as e:
        print(f"❌ Supported formats test error: {e}")

def test_word_to_text():
    """Test Word to text conversion."""
    try:
        if not os.path.exists(TEST_DOCX_PATH):
            create_test_docx()
        
        with open(TEST_DOCX_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/word-to-text", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ Word to text conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ Word to text conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ Word to text conversion test error: {e}")

def test_powerpoint_to_text():
    """Test PowerPoint to text conversion."""
    try:
        if not os.path.exists(TEST_PPTX_PATH):
            create_test_pptx()
        
        with open(TEST_PPTX_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/powerpoint-to-text", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PowerPoint to text conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PowerPoint to text conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PowerPoint to text conversion test error: {e}")

def test_pdf_to_text():
    """Test PDF to text conversion."""
    try:
        if not create_test_pdf():
            print("⚠️ PDF to text test requires PDF creation - skipping")
            return
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-text", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to text conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to text conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to text conversion test error: {e}")

def test_srt_to_text():
    """Test SRT to text conversion."""
    try:
        if not os.path.exists(TEST_SRT_PATH):
            create_test_srt()
        
        with open(TEST_SRT_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/srt-to-text", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ SRT to text conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ SRT to text conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ SRT to text conversion test error: {e}")

def test_vtt_to_text():
    """Test VTT to text conversion."""
    try:
        if not os.path.exists(TEST_VTT_PATH):
            create_test_vtt()
        
        with open(TEST_VTT_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/vtt-to-text", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ VTT to text conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ VTT to text conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ VTT to text conversion test error: {e}")

def cleanup():
    """Clean up test files."""
    test_files = [TEST_DOCX_PATH, TEST_PPTX_PATH, TEST_PDF_PATH, TEST_SRT_PATH, TEST_VTT_PATH]
    for file_path in test_files:
        if os.path.exists(file_path):
            os.remove(file_path)
            print(f"Cleaned up test file: {file_path}")

def main():
    """Run all text conversion tests."""
    print("🧪 Testing Text Conversion API")
    print("=" * 50)
    
    # Test supported formats
    print("\n1. Testing supported formats...")
    test_supported_formats()
    
    # Test Word to text
    print("\n2. Testing Word to text conversion...")
    test_word_to_text()
    
    # Test PowerPoint to text
    print("\n3. Testing PowerPoint to text conversion...")
    test_powerpoint_to_text()
    
    # Test PDF to text
    print("\n4. Testing PDF to text conversion...")
    test_pdf_to_text()
    
    # Test SRT to text
    print("\n5. Testing SRT to text conversion...")
    test_srt_to_text()
    
    # Test VTT to text
    print("\n6. Testing VTT to text conversion...")
    test_vtt_to_text()
    
    # Cleanup
    print("\n7. Cleaning up...")
    cleanup()
    
    print("\n✅ All text conversion tests completed!")

if __name__ == "__main__":
    main()
