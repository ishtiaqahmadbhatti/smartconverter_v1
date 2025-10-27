#!/usr/bin/env python3
"""
Test script for Office Documents conversion functionality.
"""

import requests
import os
import json
import tempfile

# Test configuration
BASE_URL = "http://localhost:8000/api/v1/officedocumentsconversiontools"

def create_test_docx():
    """Create test DOCX file."""
    try:
        from docx import Document
        
        doc = Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('This is a test document for Office Documents conversion.')
        doc.add_paragraph('It contains multiple paragraphs and formatting.')
        
        # Add a table
        table = doc.add_table(rows=3, cols=3)
        table.style = 'Table Grid'
        
        # Add data to table
        for i in range(3):
            for j in range(3):
                table.cell(i, j).text = f'Cell {i+1},{j+1}'
        
        doc.save('test_document.docx')
        print("Created test DOCX file: test_document.docx")
        
    except ImportError:
        print("⚠️ python-docx not available - creating dummy DOCX file")
        with open('test_document.docx', 'wb') as f:
            f.write(b"dummy docx content")

def create_test_xlsx():
    """Create test XLSX file."""
    try:
        import pandas as pd
        
        data = {
            'Product': ['Laptop', 'Mouse', 'Keyboard'],
            'Price': [999, 25, 75],
            'Quantity': [10, 50, 30],
            'Total': [9990, 1250, 2250]
        }
        
        df = pd.DataFrame(data)
        df.to_excel('test_spreadsheet.xlsx', index=False)
        print("Created test XLSX file: test_spreadsheet.xlsx")
        
    except ImportError:
        print("⚠️ pandas not available - creating dummy XLSX file")
        with open('test_spreadsheet.xlsx', 'wb') as f:
            f.write(b"dummy xlsx content")

def create_test_pptx():
    """Create test PPTX file."""
    try:
        from pptx import Presentation
        
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Test Presentation"
        subtitle.text = "This is a test PowerPoint presentation"
        
        prs.save('test_presentation.pptx')
        print("Created test PPTX file: test_presentation.pptx")
        
    except ImportError:
        print("⚠️ python-pptx not available - creating dummy PPTX file")
        with open('test_presentation.pptx', 'wb') as f:
            f.write(b"dummy pptx content")

def test_supported_formats():
    """Test getting supported formats."""
    try:
        response = requests.get(f"{BASE_URL}/supported-formats")
        if response.status_code == 200:
            print("✅ Supported formats endpoint working")
            print(f"Response: {response.json()}")
        else:
            print(f"❌ Supported formats failed: {response.status_code}")
    except Exception as e:
        print(f"❌ Supported formats error: {e}")

def test_docx_to_pdf():
    """Test DOCX to PDF conversion."""
    try:
        create_test_docx()
        
        with open('test_document.docx', 'rb') as f:
            files = {'file': ('test_document.docx', f, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')}
            response = requests.post(f"{BASE_URL}/docx-to-pdf", files=files)
        
        if response.status_code == 200:
            print("✅ DOCX to PDF conversion successful")
            with open('output.pdf', 'wb') as f:
                f.write(response.content)
            print("Saved output as output.pdf")
        else:
            print(f"❌ DOCX to PDF failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ DOCX to PDF error: {e}")

def test_xlsx_to_pdf():
    """Test XLSX to PDF conversion."""
    try:
        create_test_xlsx()
        
        with open('test_spreadsheet.xlsx', 'rb') as f:
            files = {'file': ('test_spreadsheet.xlsx', f, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')}
            response = requests.post(f"{BASE_URL}/xlsx-to-pdf", files=files)
        
        if response.status_code == 200:
            print("✅ XLSX to PDF conversion successful")
            with open('output.pdf', 'wb') as f:
                f.write(response.content)
            print("Saved output as output.pdf")
        else:
            print(f"❌ XLSX to PDF failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ XLSX to PDF error: {e}")

def test_pptx_to_pdf():
    """Test PPTX to PDF conversion."""
    try:
        create_test_pptx()
        
        with open('test_presentation.pptx', 'rb') as f:
            files = {'file': ('test_presentation.pptx', f, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
            response = requests.post(f"{BASE_URL}/pptx-to-pdf", files=files)
        
        if response.status_code == 200:
            print("✅ PPTX to PDF conversion successful")
            with open('output.pdf', 'wb') as f:
                f.write(response.content)
            print("Saved output as output.pdf")
        else:
            print(f"❌ PPTX to PDF failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PPTX to PDF error: {e}")

def test_pdf_to_docx():
    """Test PDF to DOCX conversion."""
    try:
        # Create a simple PDF first (this would normally be done by the service)
        pdf_content = b"%PDF-1.4\n1 0 obj\n<<\n/Type /Catalog\n/Pages 2 0 R\n>>\nendobj\n2 0 obj\n<<\n/Type /Pages\n/Kids [3 0 R]\n/Count 1\n>>\nendobj\n3 0 obj\n<<\n/Type /Page\n/Parent 2 0 R\n/MediaBox [0 0 612 792]\n>>\nendobj\nxref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \ntrailer\n<<\n/Size 4\n/Root 1 0 R\n>>\nstartxref\n174\n%%EOF"
        
        with open('test_input.pdf', 'wb') as f:
            f.write(pdf_content)
        
        with open('test_input.pdf', 'rb') as f:
            files = {'file': ('test_input.pdf', f, 'application/pdf')}
            response = requests.post(f"{BASE_URL}/pdf-to-docx", files=files)
        
        if response.status_code == 200:
            print("✅ PDF to DOCX conversion successful")
            with open('output.docx', 'wb') as f:
                f.write(response.content)
            print("Saved output as output.docx")
        else:
            print(f"❌ PDF to DOCX failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to DOCX error: {e}")

def test_docx_to_txt():
    """Test DOCX to TXT conversion."""
    try:
        create_test_docx()
        
        with open('test_document.docx', 'rb') as f:
            files = {'file': ('test_document.docx', f, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')}
            response = requests.post(f"{BASE_URL}/docx-to-txt", files=files)
        
        if response.status_code == 200:
            print("✅ DOCX to TXT conversion successful")
            with open('output.txt', 'wb') as f:
                f.write(response.content)
            print("Saved output as output.txt")
        else:
            print(f"❌ DOCX to TXT failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ DOCX to TXT error: {e}")

def test_xlsx_to_csv():
    """Test XLSX to CSV conversion."""
    try:
        create_test_xlsx()
        
        with open('test_spreadsheet.xlsx', 'rb') as f:
            files = {'file': ('test_spreadsheet.xlsx', f, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')}
            response = requests.post(f"{BASE_URL}/xlsx-to-csv", files=files)
        
        if response.status_code == 200:
            print("✅ XLSX to CSV conversion successful")
            with open('output.csv', 'wb') as f:
                f.write(response.content)
            print("Saved output as output.csv")
        else:
            print(f"❌ XLSX to CSV failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ XLSX to CSV error: {e}")

def cleanup():
    """Clean up test files."""
    test_files = [
        'test_document.docx', 'test_spreadsheet.xlsx', 'test_presentation.pptx',
        'test_input.pdf', 'output.pdf', 'output.docx', 'output.txt', 'output.csv'
    ]
    
    for file in test_files:
        if os.path.exists(file):
            os.remove(file)
            print(f"Cleaned up: {file}")

def main():
    """Run all tests."""
    print("🧪 Testing Office Documents Conversion API")
    print("=" * 50)
    
    test_supported_formats()
    print()
    
    test_docx_to_pdf()
    print()
    
    test_xlsx_to_pdf()
    print()
    
    test_pptx_to_pdf()
    print()
    
    test_pdf_to_docx()
    print()
    
    test_docx_to_txt()
    print()
    
    test_xlsx_to_csv()
    print()
    
    cleanup()
    print("✅ Office Documents conversion tests completed!")

if __name__ == "__main__":
    main()
