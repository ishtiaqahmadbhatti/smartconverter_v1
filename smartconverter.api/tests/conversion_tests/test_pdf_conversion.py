#!/usr/bin/env python3
"""
Test script for PDF conversion functionality.
"""

import requests
import os
import json
import tempfile

# Test configuration
BASE_URL = "http://localhost:8000/api/v1/pdfconversiontools"
TEST_PDF_PATH = "test_document.pdf"
TEST_HTML_PATH = "test_document.html"
TEST_DOCX_PATH = "test_document.docx"
TEST_PPTX_PATH = "test_presentation.pptx"
TEST_JPG_PATH = "test_image.jpg"
TEST_PNG_PATH = "test_image.png"
TEST_MD_PATH = "test_document.md"
TEST_XLSX_PATH = "test_spreadsheet.xlsx"

def create_test_files():
    """Create test files for conversion."""
    try:
        # Create a simple test PDF using reportlab
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        
        # Create PDF
        doc = SimpleDocTemplate(TEST_PDF_PATH, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        story.append(Paragraph("Test Document", styles['Title']))
        story.append(Spacer(1, 12))
        story.append(Paragraph("This is a test document for PDF conversion.", styles['Normal']))
        story.append(Paragraph("It contains multiple paragraphs and formatting.", styles['Normal']))
        story.append(Spacer(1, 12))
        story.append(Paragraph("Table Data:", styles['Heading2']))
        story.append(Paragraph("Name, Age, City", styles['Normal']))
        story.append(Paragraph("John, 25, New York", styles['Normal']))
        story.append(Paragraph("Jane, 30, London", styles['Normal']))
        
        doc.build(story)
        
        # Create HTML
        html_content = """
        <!DOCTYPE html>
        <html>
        <head><title>Test Document</title></head>
        <body>
            <h1>Test Document</h1>
            <p>This is a test HTML document for conversion.</p>
            <p>It contains multiple paragraphs and formatting.</p>
            <table>
                <tr><th>Name</th><th>Age</th><th>City</th></tr>
                <tr><td>John</td><td>25</td><td>New York</td></tr>
                <tr><td>Jane</td><td>30</td><td>London</td></tr>
            </table>
        </body>
        </html>
        """
        with open(TEST_HTML_PATH, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        # Create Word document
        from docx import Document
        doc = Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('This is a test Word document for conversion.')
        doc.add_paragraph('It contains multiple paragraphs and formatting.')
        doc.add_heading('Table Data:', level=1)
        table = doc.add_table(rows=1, cols=3)
        table.cell(0, 0).text = 'Name'
        table.cell(0, 1).text = 'Age'
        table.cell(0, 2).text = 'City'
        row = table.add_row()
        row.cells[0].text = 'John'
        row.cells[1].text = '25'
        row.cells[2].text = 'New York'
        doc.save(TEST_DOCX_PATH)
        
        # Create PowerPoint
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Presentation"
        slide.placeholders[1].text = "This is a test PowerPoint presentation for conversion."
        prs.save(TEST_PPTX_PATH)
        
        # Create images
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new('RGB', (400, 300), color='white')
        draw = ImageDraw.Draw(img)
        draw.text((50, 50), "Test Image", fill='black')
        draw.text((50, 100), "This is a test image for conversion.", fill='black')
        img.save(TEST_JPG_PATH, 'JPEG')
        img.save(TEST_PNG_PATH, 'PNG')
        
        # Create Markdown
        md_content = """# Test Document

This is a test Markdown document for conversion.

It contains multiple paragraphs and formatting.

## Table Data

| Name | Age | City |
|------|-----|------|
| John | 25  | New York |
| Jane | 30  | London |
"""
        with open(TEST_MD_PATH, 'w', encoding='utf-8') as f:
            f.write(md_content)
        
        # Create Excel
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Test Sheet"
        ws['A1'] = 'Name'
        ws['B1'] = 'Age'
        ws['C1'] = 'City'
        ws['A2'] = 'John'
        ws['B2'] = 25
        ws['C2'] = 'New York'
        ws['A3'] = 'Jane'
        ws['B3'] = 30
        ws['C3'] = 'London'
        wb.save(TEST_XLSX_PATH)
        
        print(f"Created test files: {TEST_PDF_PATH}, {TEST_HTML_PATH}, {TEST_DOCX_PATH}, {TEST_PPTX_PATH}, {TEST_JPG_PATH}, {TEST_PNG_PATH}, {TEST_MD_PATH}, {TEST_XLSX_PATH}")
        
    except ImportError as e:
        print(f"⚠️ Required libraries not available: {e}")
        print("Creating dummy files...")
        # Create dummy files
        for file_path in [TEST_PDF_PATH, TEST_HTML_PATH, TEST_DOCX_PATH, TEST_PPTX_PATH, TEST_JPG_PATH, TEST_PNG_PATH, TEST_MD_PATH, TEST_XLSX_PATH]:
            with open(file_path, 'wb') as f:
                f.write(b"dummy content")

def test_supported_formats():
    """Test getting supported formats."""
    try:
        response = requests.get(f"{BASE_URL}/supported-formats")
        if response.status_code == 200:
            data = response.json()
            print("✅ Supported formats test passed")
            print(f"Supported formats: {json.dumps(data, indent=2)}")
        else:
            print(f"❌ Supported formats test failed: {response.status_code}")
    except Exception as e:
        print(f"❌ Supported formats test error: {e}")

def test_pdf_to_json():
    """Test AI: Convert PDF to JSON."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-json", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ AI: PDF to JSON conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ AI: PDF to JSON conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ AI: PDF to JSON conversion test error: {e}")

def test_pdf_to_markdown():
    """Test AI: Convert PDF to Markdown."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-markdown", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ AI: PDF to Markdown conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ AI: PDF to Markdown conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ AI: PDF to Markdown conversion test error: {e}")

def test_pdf_to_csv():
    """Test AI: Convert PDF to CSV."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-csv", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ AI: PDF to CSV conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ AI: PDF to CSV conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ AI: PDF to CSV conversion test error: {e}")

def test_pdf_to_excel():
    """Test AI: Convert PDF to Excel."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-excel", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ AI: PDF to Excel conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ AI: PDF to Excel conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ AI: PDF to Excel conversion test error: {e}")

def test_html_to_pdf():
    """Test Convert HTML to PDF."""
    try:
        if not os.path.exists(TEST_HTML_PATH):
            create_test_files()
        
        with open(TEST_HTML_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/html-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ HTML to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ HTML to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ HTML to PDF conversion test error: {e}")

def test_word_to_pdf():
    """Test Convert Word to PDF."""
    try:
        if not os.path.exists(TEST_DOCX_PATH):
            create_test_files()
        
        with open(TEST_DOCX_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/word-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ Word to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ Word to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ Word to PDF conversion test error: {e}")

def test_powerpoint_to_pdf():
    """Test Convert PowerPoint to PDF."""
    try:
        if not os.path.exists(TEST_PPTX_PATH):
            create_test_files()
        
        with open(TEST_PPTX_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/powerpoint-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PowerPoint to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PowerPoint to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PowerPoint to PDF conversion test error: {e}")

def test_jpg_to_pdf():
    """Test Convert JPG to PDF."""
    try:
        if not os.path.exists(TEST_JPG_PATH):
            create_test_files()
        
        with open(TEST_JPG_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/jpg-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ JPG to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ JPG to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ JPG to PDF conversion test error: {e}")

def test_png_to_pdf():
    """Test Convert PNG to PDF."""
    try:
        if not os.path.exists(TEST_PNG_PATH):
            create_test_files()
        
        with open(TEST_PNG_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/png-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PNG to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PNG to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PNG to PDF conversion test error: {e}")

def test_markdown_to_pdf():
    """Test Convert Markdown to PDF."""
    try:
        if not os.path.exists(TEST_MD_PATH):
            create_test_files()
        
        with open(TEST_MD_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/markdown-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ Markdown to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ Markdown to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ Markdown to PDF conversion test error: {e}")

def test_excel_to_pdf():
    """Test Convert Excel to PDF."""
    try:
        if not os.path.exists(TEST_XLSX_PATH):
            create_test_files()
        
        with open(TEST_XLSX_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/excel-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ Excel to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ Excel to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ Excel to PDF conversion test error: {e}")

def test_pdf_to_jpg():
    """Test Convert PDF to JPG."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-jpg", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to JPG conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to JPG conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to JPG conversion test error: {e}")

def test_pdf_to_png():
    """Test Convert PDF to PNG."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-png", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to PNG conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to PNG conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to PNG conversion test error: {e}")

def test_pdf_to_tiff():
    """Test Convert PDF to TIFF."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-tiff", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to TIFF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to TIFF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to TIFF conversion test error: {e}")

def test_pdf_to_svg():
    """Test Convert PDF to SVG."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-svg", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to SVG conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to SVG conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to SVG conversion test error: {e}")

def test_pdf_to_html():
    """Test Convert PDF to HTML."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-html", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to HTML conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to HTML conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to HTML conversion test error: {e}")

def test_pdf_to_text():
    """Test Convert PDF to Text."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-text", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to Text conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to Text conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to Text conversion test error: {e}")

def test_pdf_to_csv_extract():
    """Test Convert PDF to CSV (Extract)."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-csv-extract", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to CSV (Extract) conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to CSV (Extract) conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to CSV (Extract) conversion test error: {e}")

def test_pdf_to_excel_extract():
    """Test Convert PDF to Excel (Extract)."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-excel-extract", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to Excel (Extract) conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to Excel (Extract) conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to Excel (Extract) conversion test error: {e}")

def test_pdf_to_word_extract():
    """Test Convert PDF to Word (Extract)."""
    try:
        if not os.path.exists(TEST_PDF_PATH):
            create_test_files()
        
        with open(TEST_PDF_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/pdf-to-word-extract", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ PDF to Word (Extract) conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ PDF to Word (Extract) conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ PDF to Word (Extract) conversion test error: {e}")

def test_oxps_to_pdf():
    """Test Convert OXPS to PDF."""
    try:
        # Create dummy OXPS file
        oxps_path = "test_document.oxps"
        with open(oxps_path, 'wb') as f:
            f.write(b"dummy oxps content")
        
        with open(oxps_path, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/oxps-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ OXPS to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ OXPS to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ OXPS to PDF conversion test error: {e}")

def test_excel_to_xps():
    """Test Convert Excel to XPS."""
    try:
        if not os.path.exists(TEST_XLSX_PATH):
            create_test_files()
        
        with open(TEST_XLSX_PATH, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/excel-to-xps", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ Excel to XPS conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ Excel to XPS conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ Excel to XPS conversion test error: {e}")

def test_ods_to_pdf():
    """Test Convert OpenOffice Calc ODS to PDF."""
    try:
        # Create dummy ODS file
        ods_path = "test_spreadsheet.ods"
        with open(ods_path, 'wb') as f:
            f.write(b"dummy ods content")
        
        with open(ods_path, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{BASE_URL}/ods-to-pdf", files=files)
        
        if response.status_code == 200:
            result = response.json()
            print("✅ ODS to PDF conversion test passed")
            print(f"Result: {json.dumps(result, indent=2)}")
        else:
            print(f"❌ ODS to PDF conversion test failed: {response.status_code}")
            print(f"Error: {response.text}")
    except Exception as e:
        print(f"❌ ODS to PDF conversion test error: {e}")

def cleanup():
    """Clean up test files."""
    test_files = [
        TEST_PDF_PATH, TEST_HTML_PATH, TEST_DOCX_PATH, TEST_PPTX_PATH, 
        TEST_JPG_PATH, TEST_PNG_PATH, TEST_MD_PATH, TEST_XLSX_PATH,
        "test_document.oxps", "test_spreadsheet.ods"
    ]
    for file_path in test_files:
        if os.path.exists(file_path):
            os.remove(file_path)
            print(f"Cleaned up test file: {file_path}")

def main():
    """Run all PDF conversion tests."""
    print("🧪 Testing PDF Conversion API")
    print("=" * 50)
    
    # Test supported formats
    print("\n1. Testing supported formats...")
    test_supported_formats()
    
    # AI-Powered Conversions
    print("\n2. Testing AI: PDF to JSON conversion...")
    test_pdf_to_json()
    
    print("\n3. Testing AI: PDF to Markdown conversion...")
    test_pdf_to_markdown()
    
    print("\n4. Testing AI: PDF to CSV conversion...")
    test_pdf_to_csv()
    
    print("\n5. Testing AI: PDF to Excel conversion...")
    test_pdf_to_excel()
    
    # Document to PDF Conversions
    print("\n6. Testing HTML to PDF conversion...")
    test_html_to_pdf()
    
    print("\n7. Testing Word to PDF conversion...")
    test_word_to_pdf()
    
    print("\n8. Testing PowerPoint to PDF conversion...")
    test_powerpoint_to_pdf()
    
    print("\n9. Testing OXPS to PDF conversion...")
    test_oxps_to_pdf()
    
    print("\n10. Testing JPG to PDF conversion...")
    test_jpg_to_pdf()
    
    print("\n11. Testing PNG to PDF conversion...")
    test_png_to_pdf()
    
    print("\n12. Testing Markdown to PDF conversion...")
    test_markdown_to_pdf()
    
    print("\n13. Testing Excel to PDF conversion...")
    test_excel_to_pdf()
    
    print("\n14. Testing Excel to XPS conversion...")
    test_excel_to_xps()
    
    print("\n15. Testing ODS to PDF conversion...")
    test_ods_to_pdf()
    
    # PDF to Other Format Conversions
    print("\n16. Testing PDF to CSV (Extract) conversion...")
    test_pdf_to_csv_extract()
    
    print("\n17. Testing PDF to Excel (Extract) conversion...")
    test_pdf_to_excel_extract()
    
    print("\n18. Testing PDF to Word (Extract) conversion...")
    test_pdf_to_word_extract()
    
    print("\n19. Testing PDF to JPG conversion...")
    test_pdf_to_jpg()
    
    print("\n20. Testing PDF to PNG conversion...")
    test_pdf_to_png()
    
    print("\n21. Testing PDF to TIFF conversion...")
    test_pdf_to_tiff()
    
    print("\n22. Testing PDF to SVG conversion...")
    test_pdf_to_svg()
    
    print("\n23. Testing PDF to HTML conversion...")
    test_pdf_to_html()
    
    print("\n24. Testing PDF to Text conversion...")
    test_pdf_to_text()
    
    # Cleanup
    print("\n25. Cleaning up...")
    cleanup()
    
    print("\n✅ All PDF conversion tests completed!")

if __name__ == "__main__":
    main()